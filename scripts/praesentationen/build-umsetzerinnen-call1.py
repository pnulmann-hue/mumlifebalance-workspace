"""
Die Umsetzerinnen - Call 1: Die 5 Lecks im Online-Business
20-Folien-PPTX mit Patricia-Brand + Sprechnotizen.

Build:
    python scripts/praesentationen/build-umsetzerinnen-call1.py

Output:
    outputs/praesentationen/die-umsetzerinnen-call1-5-lecks.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# === Brand ===
CREME = RGBColor(0xF1, 0xEC, 0xDD)
DUNKELBLAU = RGBColor(0x29, 0x55, 0x6D)
PETROL = RGBColor(0x12, 0x82, 0x8C)
ORANGE = RGBColor(0xDC, 0x82, 0x2E)
TEXT_DARK = RGBColor(0x0C, 0x1C, 0x30)
WEISS = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEAD = "Cambria"   # Philosopher-Fallback (jeder Win-PC hat Cambria)
FONT_BODY = "Calibri"   # Source-Sans-3-Fallback

# === Setup ===
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]


# === Helper ===
def add_bg(slide, color=CREME):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_top_bar(slide, color=PETROL, height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_bottom_bar(slide, color=PETROL, height=Inches(0.06)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SH - height, SW, height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    return bar


def add_text(
    slide, text, left, top, width, height,
    *, font=FONT_BODY, size=18, color=TEXT_DARK, bold=False, italic=False,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.2,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_eyebrow(slide, text, left, top):
    add_text(
        slide, text.upper(), left, top, Inches(8), Inches(0.4),
        font=FONT_BODY, size=14, color=DUNKELBLAU, bold=True,
        align=PP_ALIGN.LEFT,
    )


def add_footer(slide, page_num, total=20):
    add_text(
        slide, f"Die Umsetzerinnen  ·  Call 1  ·  {page_num}/{total}",
        Inches(0.6), SH - Inches(0.45), Inches(12.1), Inches(0.3),
        font=FONT_BODY, size=10, color=DUNKELBLAU, align=PP_ALIGN.LEFT,
    )


def add_orange_divider(slide, left, top, width=Inches(0.7)):
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05))
    d.line.fill.background()
    d.fill.solid()
    d.fill.fore_color.rgb = ORANGE


def add_big_number(slide, n, left, top):
    """Riesige Petrol-Zahl als Hintergrund-Element."""
    add_text(
        slide, str(n), left, top, Inches(4), Inches(4),
        font=FONT_HEAD, size=380, color=PETROL, bold=True,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=0.9,
    )


def set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    p = notes.paragraphs[0]
    p.text = text


# === Layout-Templates ===
def slide_title_hero(eyebrow, title, subtitle, footer_meta, notes):
    s = prs.slides.add_slide(BLANK)
    # Dunkelblauer Vollflächen-Hintergrund fuer Titelfolie
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = DUNKELBLAU

    # Petrol Akzentstreifen seitlich
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SH)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE

    add_text(
        s, eyebrow.upper(), Inches(1.0), Inches(1.4),
        Inches(11), Inches(0.5),
        font=FONT_BODY, size=18, color=CREME, bold=True,
        align=PP_ALIGN.LEFT,
    )
    add_text(
        s, title, Inches(1.0), Inches(2.1),
        Inches(11), Inches(2.5),
        font=FONT_HEAD, size=88, color=CREME, bold=True,
        align=PP_ALIGN.LEFT, line_spacing=1.0,
    )
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.8), Inches(0.9), Inches(0.06))
    div.line.fill.background()
    div.fill.solid()
    div.fill.fore_color.rgb = ORANGE
    add_text(
        s, subtitle, Inches(1.0), Inches(5.05),
        Inches(11), Inches(1.0),
        font=FONT_HEAD, size=30, color=CREME, italic=True,
        align=PP_ALIGN.LEFT,
    )
    add_text(
        s, footer_meta, Inches(1.0), Inches(6.6),
        Inches(11), Inches(0.4),
        font=FONT_BODY, size=14, color=CREME, align=PP_ALIGN.LEFT,
    )
    set_notes(s, notes)
    return s


def slide_section_cover(number, title, subtitle, notes, page_num):
    """Volle Cover-Folie pro Leck mit grosser Zahl."""
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_top_bar(s)
    add_bottom_bar(s)

    # Linke Seite: grosse Zahl
    add_text(
        s, str(number), Inches(0.8), Inches(1.3),
        Inches(5), Inches(5),
        font=FONT_HEAD, size=380, color=PETROL, bold=True,
        align=PP_ALIGN.LEFT, line_spacing=0.85,
    )

    # Rechte Seite
    add_eyebrow(s, f"Leck {number} von 5", Inches(6.5), Inches(2.2))
    add_text(
        s, title, Inches(6.5), Inches(2.6),
        Inches(6.3), Inches(2.0),
        font=FONT_HEAD, size=58, color=DUNKELBLAU, bold=True,
        align=PP_ALIGN.LEFT, line_spacing=1.05,
    )
    add_orange_divider(s, Inches(6.5), Inches(4.7))
    add_text(
        s, subtitle, Inches(6.5), Inches(4.95),
        Inches(6.3), Inches(1.5),
        font=FONT_HEAD, size=22, color=TEXT_DARK, italic=True,
        align=PP_ALIGN.LEFT, line_spacing=1.3,
    )
    add_footer(s, page_num)
    set_notes(s, notes)
    return s


def slide_content(eyebrow, title, body_lines, notes, page_num,
                  body_size=22, list_style="dot"):
    """Standard-Content-Folie: Eyebrow + Title + Liste/Body."""
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_top_bar(s)
    add_bottom_bar(s)

    add_eyebrow(s, eyebrow, Inches(0.8), Inches(0.6))
    add_text(
        s, title, Inches(0.8), Inches(1.05),
        Inches(11.7), Inches(1.4),
        font=FONT_HEAD, size=42, color=DUNKELBLAU, bold=True,
        align=PP_ALIGN.LEFT, line_spacing=1.05,
    )
    add_orange_divider(s, Inches(0.8), Inches(2.55))

    # Body
    body_text = ""
    for line in body_lines:
        if list_style == "dot":
            body_text += f"•  {line}\n"
        elif list_style == "arrow":
            body_text += f"→  {line}\n"
        else:
            body_text += f"{line}\n"
    body_text = body_text.rstrip()

    add_text(
        s, body_text, Inches(0.8), Inches(2.95),
        Inches(11.7), Inches(4.1),
        font=FONT_BODY, size=body_size, color=TEXT_DARK,
        align=PP_ALIGN.LEFT, line_spacing=1.5,
    )
    add_footer(s, page_num)
    set_notes(s, notes)
    return s


def slide_two_col(eyebrow, title, left_head, left_lines, right_head, right_lines,
                  notes, page_num, left_color=ORANGE, right_color=PETROL):
    """Zwei-Spalten-Vergleich (z.B. Symptom links / Fix rechts)."""
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_top_bar(s)
    add_bottom_bar(s)

    add_eyebrow(s, eyebrow, Inches(0.8), Inches(0.6))
    add_text(
        s, title, Inches(0.8), Inches(1.05),
        Inches(11.7), Inches(1.4),
        font=FONT_HEAD, size=38, color=DUNKELBLAU, bold=True,
        align=PP_ALIGN.LEFT, line_spacing=1.05,
    )
    add_orange_divider(s, Inches(0.8), Inches(2.45))

    col_w = Inches(5.6)
    col_top = Inches(2.85)
    col_h = Inches(4.0)

    # Linke Spalte
    add_text(
        s, left_head.upper(), Inches(0.8), col_top,
        col_w, Inches(0.5),
        font=FONT_BODY, size=16, color=left_color, bold=True,
        align=PP_ALIGN.LEFT,
    )
    body_l = "\n".join(f"•  {x}" for x in left_lines)
    add_text(
        s, body_l, Inches(0.8), col_top + Inches(0.55),
        col_w, col_h,
        font=FONT_BODY, size=18, color=TEXT_DARK, line_spacing=1.5,
    )

    # Rechte Spalte
    add_text(
        s, right_head.upper(), Inches(7.0), col_top,
        col_w, Inches(0.5),
        font=FONT_BODY, size=16, color=right_color, bold=True,
        align=PP_ALIGN.LEFT,
    )
    body_r = "\n".join(f"→  {x}" for x in right_lines)
    add_text(
        s, body_r, Inches(0.8) + Inches(6.2), col_top + Inches(0.55),
        col_w, col_h,
        font=FONT_BODY, size=18, color=TEXT_DARK, line_spacing=1.5,
    )

    add_footer(s, page_num)
    set_notes(s, notes)
    return s


def slide_quote_break(quote, attribution, notes, page_num):
    """Volle dunkelblaue Folie mit Zitat - visueller Reset zwischendurch."""
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = DUNKELBLAU

    add_text(
        s, "“", Inches(1.0), Inches(1.0),
        Inches(2), Inches(2),
        font=FONT_HEAD, size=200, color=ORANGE, bold=True,
        align=PP_ALIGN.LEFT,
    )
    add_text(
        s, quote, Inches(1.5), Inches(2.6),
        Inches(10.3), Inches(3),
        font=FONT_HEAD, size=44, color=CREME, italic=True, bold=False,
        align=PP_ALIGN.LEFT, line_spacing=1.25,
    )
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.8), Inches(0.7), Inches(0.05))
    div.line.fill.background()
    div.fill.solid()
    div.fill.fore_color.rgb = ORANGE
    add_text(
        s, f"- {attribution}", Inches(1.5), Inches(6.0),
        Inches(10), Inches(0.5),
        font=FONT_BODY, size=16, color=CREME, align=PP_ALIGN.LEFT,
    )
    add_text(
        s, f"Die Umsetzerinnen  ·  Call 1  ·  {page_num}/20",
        Inches(0.6), SH - Inches(0.45), Inches(12.1), Inches(0.3),
        font=FONT_BODY, size=10, color=CREME, align=PP_ALIGN.LEFT,
    )
    set_notes(s, notes)
    return s


# ============================================================================
# === FOLIEN ===
# ============================================================================

# Folie 1 - Titelfolie
slide_title_hero(
    eyebrow="Mum Life Balance  ·  Patricia Ulmann",
    title="Die\nUmsetzerinnen",
    subtitle="Nicht uberlegen. Nicht warten. Umsetzen.",
    footer_meta="Call 1  ·  30. April 2026  ·  09:00 Uhr",
    notes=(
        "BEGRUESSUNG (1 Min):\n"
        "- Hallo zusammen, schoen dass ihr da seid.\n"
        "- Heute starten wir mit Call 1: Die 5 Lecks im Online-Business.\n"
        "- Was ihr heute lernt, koennt ihr ab morgen umsetzen.\n"
        "- Ihr seid 6 Frauen - kleine Runde, alle bekommen Raum."
    ),
)

# Folie 2 - Was die Umsetzerinnen sind
slide_content(
    eyebrow="Willkommen",
    title="Was Umsetzerinnen anders machen",
    body_lines=[
        "Sie konsumieren keine Tipps mehr - sie testen sie.",
        "Sie warten nicht auf Perfektion - sie veroeffentlichen B-Versionen.",
        "Sie hoeren auf, nach dem naechsten Kurs zu suchen - sie machen den, den sie schon haben.",
        "Sie messen Ergebnisse - nicht Aufwand.",
    ],
    notes=(
        "WORUM ES IN DIESEM PROGRAMM GEHT (3 Min):\n"
        "- Es gibt einen Unterschied zwischen Lernen und Umsetzen.\n"
        "- Die meisten haengen in der Lern-Schleife: Webinar, Podcast, naechstes Webinar.\n"
        "- Hier brechen wir das. Jeder Call hat ein konkretes Umsetzungs-Ergebnis.\n"
        "- Frage in die Runde: Wer kennt das Gefuehl, viel zu wissen aber wenig zu posten?"
    ),
    page_num=2,
    body_size=24,
)

# Folie 3 - Agenda
slide_content(
    eyebrow="Heute",
    title="Die 5 Lecks - was dich Umsatz kostet",
    body_lines=[
        "Bio-Leck  -  Dein Profil verkauft nicht fuer dich.",
        "Reichweite-Leck  -  Du wirst gar nicht erst gesehen.",
        "Lead-Leck  -  Besucher werden nicht zu Leads.",
        "Welcome-Leck  -  Leads werden nicht zu Kundinnen.",
        "Verkaufs-Leck  -  Kundinnen kaufen einmal - nie wieder.",
    ],
    notes=(
        "AGENDA (1 Min):\n"
        "- Wir gehen heute durch 5 Stellen, an denen euer Business 'leckt'.\n"
        "- Jedes Leck ist ein Punkt, an dem Reichweite, Leads oder Geld verloren gehen.\n"
        "- Pro Leck: Was es ist, woran ihr es erkennt, der Fix.\n"
        "- Am Ende sagt ihr mir: welches Leck blutet bei EUCH am staerksten."
    ),
    page_num=3,
    body_size=20,
    list_style="dot",
)

# Folie 4 - Eimer-Metapher
slide_content(
    eyebrow="Die Metapher",
    title="Stell dir dein Business als Eimer vor.",
    body_lines=[
        "Oben kommt Wasser rein - das ist Reichweite (Reels, Posts, Stories).",
        "Unten soll Umsatz rauskommen - das sind Kaeufe.",
        "Aber dazwischen sind 5 Lecks, durch die staendig Wasser entweicht.",
        "Mehr Reichweite produzieren = mehr Wasser oben reinkippen.",
        "Lecks abdichten = der Eimer haelt endlich, was du reinkippst.",
    ],
    notes=(
        "EIMER-METAPHER (3 Min):\n"
        "- Das ist die wichtigste Folie heute. Macht Screenshots.\n"
        "- Die meisten von euch arbeiten am Wasserhahn - mehr Reels, mehr Posts.\n"
        "- Aber wenn der Eimer leckt, ist das wie Wasser ins Sieb kippen.\n"
        "- Die Reihenfolge ist: ZUERST Lecks abdichten, DANN auf Reichweite gehen.\n"
        "- Beispiel: Wenn dein Profil nicht verkauft, bringen 10x mehr Follower 10x mehr Frust."
    ),
    page_num=4,
    body_size=20,
    list_style="arrow",
)

# === LECK 1 - BIO ===
slide_section_cover(
    number=1,
    title="Das Bio-Leck",
    subtitle="Dein Profil ist eine Visitenkarte - keine Verkaufsflaeche.",
    notes=(
        "EINSTIEG LECK 1 (1 Min):\n"
        "- Das ist das Leck, das ich am haeufigsten sehe.\n"
        "- 90% der Profile, die ich pruefe, haben hier 3-5 Fehler.\n"
        "- Erinnerung: Eure Bio ist 150 Zeichen, die entscheiden ob jemand bleibt oder geht.\n"
        "- Die meisten verschenken hier 80% ihrer Conversion."
    ),
    page_num=5,
)

slide_two_col(
    eyebrow="Leck 1 - Diagnose",
    title="So merkst du, dass dein Bio leckt",
    left_head="Symptome",
    left_lines=[
        "Leute folgen dir und schauen nie wieder rein.",
        "Niemand schreibt dir, obwohl Reichweite da ist.",
        "Du musst erklaeren, was du machst, wenn jemand fragt.",
        "Du wechselst die Bio alle 4 Wochen - und nichts aendert sich.",
    ],
    right_head="Das Bio-Leck im Klartext",
    right_lines=[
        "Dein Bio sagt WAS du bist (\"Mama, Coach, Mentorin\").",
        "Nicht: WAS bekommt sie - bis WANN - durch WEN.",
        "Kein einziger Klick-Pfad fuer Anfrage / 0EUR-Produkt.",
    ],
    notes=(
        "DIAGNOSE BIO-LECK (4 Min):\n"
        "- Lasst uns das ehrlich machen. Geht in eine eigene Bio rein.\n"
        "- Frage 1: Sieht man in 3 Sekunden, fuer WEN du da bist?\n"
        "- Frage 2: Sieht man, was sie konkret von dir bekommen?\n"
        "- Frage 3: Gibt es einen klaren naechsten Schritt - DM, Link, Freebie?\n"
        "- Wenn nur einer dieser Punkte unklar ist: Bio leckt."
    ),
    page_num=6,
)

slide_content(
    eyebrow="Leck 1 - Der Fix",
    title="Die 3-Zeilen-Bio, die verkauft",
    body_lines=[
        "Zeile 1: Fuer WEN bist du da? (Zielgruppe + Situation)",
        "Zeile 2: WAS bekommen sie und bis WANN? (Versprechen + Timeline)",
        "Zeile 3: WIE gehts weiter? (1 klarer CTA - DM, Link, Freebie)",
        "Beispiel Patricia: 'Mama, die online wachsen will - ohne 60h die Woche. Free Bio-Check unten -> Klarheit in 24h.'",
    ],
    notes=(
        "DER FIX (5 Min):\n"
        "- Diese Vorlage gilt fuer beide Profile (Mentoring + doTERRA).\n"
        "- Schreibt eure Bio JETZT auf, wir lesen 1-2 vor und feilen live.\n"
        "- Hausaufgabe: Bis zum naechsten Call neue Bio live + 1 Woche messen.\n"
        "- Tipp: 0EUR-Bio-Check ist mein Freebie - wer will, kann ich euch durchgeben."
    ),
    page_num=7,
    body_size=19,
    list_style="dot",
)

# === LECK 2 - REICHWEITE ===
slide_section_cover(
    number=2,
    title="Das Reichweite-Leck",
    subtitle="Du postest - aber niemand sieht dich.",
    notes=(
        "EINSTIEG LECK 2 (1 Min):\n"
        "- Reichweite ist nicht 'der Algorithmus mag mich nicht'.\n"
        "- Reichweite ist die Konsequenz aus 3 Stellschrauben.\n"
        "- Wir gucken jetzt welche."
    ),
    page_num=8,
)

slide_two_col(
    eyebrow="Leck 2 - Diagnose",
    title="So merkst du, dass deine Reichweite leckt",
    left_head="Symptome",
    left_lines=[
        "Reels unter 500 Views obwohl du 1000+ Follower hast.",
        "Karussells nur von gleichen 5 Leuten geliked.",
        "Nullsaves, null Shares, null DMs.",
        "Engagement-Rate unter 1%.",
    ],
    right_head="Die echten Ursachen",
    right_lines=[
        "Hook in den ersten 3 Sek nicht stark genug.",
        "Keine emotionale Identifikation - reines Tipps-Posting.",
        "Posting-Zeit passt nicht zur Zielgruppe.",
    ],
    notes=(
        "DIAGNOSE REICHWEITE-LECK (4 Min):\n"
        "- Wenn ihr Reichweite-Probleme habt, ist es fast nie der Algorithmus.\n"
        "- Es sind 3 Sachen: Hook, Identifikation, Timing.\n"
        "- Hook = die ersten 3 Sekunden / die ersten 6 Worte.\n"
        "- Identifikation = sieht sich die Leserin in deinem Post wieder?\n"
        "- Timing = postest du, wenn deine Zielgruppe wirklich am Handy ist?"
    ),
    page_num=9,
)

slide_content(
    eyebrow="Leck 2 - Der Fix",
    title="Die Hook-Reichweite-Formel",
    body_lines=[
        "Hook-Test: Wuerdest du selber stoppen, wenn du den Post siehst?",
        "Identifikations-Test: Steht 'ich' oder 'du' in den ersten 6 Worten?",
        "Timing-Test: Postest du wann SIE Zeit hat - nicht wann DU Zeit hast?",
        "Ein Test pro Woche aendern. Nicht alles auf einmal.",
    ],
    notes=(
        "DER FIX (4 Min):\n"
        "- Wichtig: nicht alles gleichzeitig aendern. EIN Hebel pro Woche.\n"
        "- Woche 1: Nur Hooks ueben. Jeden Reel mit anderem Hook.\n"
        "- Woche 2: Identifikation in den ersten 2 Saetzen.\n"
        "- Woche 3: Posting-Zeit testen.\n"
        "- Tracking-Tabelle: Views + Saves + Shares + DMs - 1x pro Woche reinschauen.\n"
        "- Mehr dazu in Call 3: 'Stories die nicht nur Likes bringen'."
    ),
    page_num=10,
    body_size=20,
    list_style="arrow",
)

# === ZWISCHEN-ZITAT ===
slide_quote_break(
    quote="Reichweite ist nicht der Job des Algorithmus.\nEs ist der Job deiner Hook.",
    attribution="Patricia",
    notes=(
        "PAUSE-MOMENT (1 Min):\n"
        "- Kurze Atempause. Macht euch eine Notiz auf diesem Satz.\n"
        "- Wer fuehlt sich gerade gemeint? Hand hoch / Reaction senden."
    ),
    page_num=11,
)

# === LECK 3 - LEAD ===
slide_section_cover(
    number=3,
    title="Das Lead-Leck",
    subtitle="Sie schauen sich um - aber gehen ohne Email-Adresse wieder.",
    notes=(
        "EINSTIEG LECK 3 (1 Min):\n"
        "- Reichweite + gutes Profil = Besucher.\n"
        "- Aber Besucher zahlen dir keine Miete.\n"
        "- Ohne Lead-Magnet ist alles, was du tust, ein Strohfeuer."
    ),
    page_num=12,
)

slide_two_col(
    eyebrow="Leck 3 - Diagnose und Fix",
    title="Du brauchst genau EIN funktionierendes Freebie",
    left_head="So merkst du, es leckt",
    left_lines=[
        "Du hast 0-2 neue Email-Eintraege pro Woche.",
        "Du hast keine Email-Liste - oder eine, die nichts macht.",
        "Du hast 3 Freebies, aber keins zieht richtig.",
        "Du verschenkst Wissen in Posts - aber sammelst keine Adressen.",
    ],
    right_head="Der Fix",
    right_lines=[
        "EIN Freebie, das EINEN spezifischen Painpoint loest.",
        "Ein klarer CTA in jedem Reel / jeder Story / Bio.",
        "Auto-Mail-1 binnen 5 Min nach Eintragung.",
    ],
    notes=(
        "DIAGNOSE + FIX LEAD-LECK (5 Min):\n"
        "- Beispiele aus meinem Business:\n"
        "  - Bio-Check (Mentoring): Tausch echte Bio gegen Feedback.\n"
        "  - Workbook 'Von 0 auf echt' (Mentoring): Klarheit in 7 Tagen.\n"
        "  - Energie-Kur (doTERRA): 14 Tage konkret + Produktrabatt.\n"
        "- Aufgabe: Welcher Painpoint, gegen den deine Zielgruppe HEUTE googelt?\n"
        "- Das ist dein Freebie. Nicht 'alles ueber Persoenlichkeitsentwicklung'."
    ),
    page_num=13,
)

# === LECK 4 - WELCOME ===
slide_section_cover(
    number=4,
    title="Das Welcome-Leck",
    subtitle="Lead landet auf der Liste - und hoert nie wieder von dir.",
    notes=(
        "EINSTIEG LECK 4 (1 Min):\n"
        "- Das ist das Leck, das mich am laengsten gekostet hat.\n"
        "- Ich hatte Eintraege - aber keine Welcome-Mails.\n"
        "- Resultat: 80% der Leads waren nach 4 Wochen kalt."
    ),
    page_num=14,
)

slide_two_col(
    eyebrow="Leck 4 - Diagnose und Fix",
    title="Die 5-Mail-Welcome-Sequenz",
    left_head="So merkst du, es leckt",
    left_lines=[
        "Du hast 200+ Email-Adressen, aber nur 1-2 Sales pro Monat.",
        "Du verschickst nur Newsletter, keine Sequenz.",
        "Oeffnungsrate unter 25%.",
        "Niemand antwortet auf deine Mails.",
    ],
    right_head="Der Fix - 5 Mails, 7 Tage",
    right_lines=[
        "Tag 0: Freebie + persoenliche Begruessung.",
        "Tag 1: Deine Story - wer du bist und WIESO.",
        "Tag 3: Ein konkretes Resultat einer Kundin.",
        "Tag 5: Dein Angebot - direkt benannt.",
        "Tag 7: Die Tuer geht zu - klare Frist.",
    ],
    notes=(
        "DIAGNOSE + FIX WELCOME-LECK (5 Min):\n"
        "- Diese Sequenz ist das Herz eures Funnels.\n"
        "- Wir bauen die in Call 2 KOMPLETT zusammen - mit Vorlagen.\n"
        "- Heute nur die Struktur, damit ihr versteht was fehlt.\n"
        "- Wichtig: schreiben in DU-Form, eine Person. Keine 'Liebe Community'."
    ),
    page_num=15,
)

# === LECK 5 - VERKAUF ===
slide_section_cover(
    number=5,
    title="Das Verkaufs-Leck",
    subtitle="Kundinnen kaufen einmal - dann nie wieder.",
    notes=(
        "EINSTIEG LECK 5 (1 Min):\n"
        "- Das letzte Leck und das tueckischste.\n"
        "- Ihr habt eine Kundin gewonnen - super.\n"
        "- Aber wenn sie nur 1x kauft, habt ihr 80% des moeglichen Wertes verschenkt."
    ),
    page_num=16,
)

slide_two_col(
    eyebrow="Leck 5 - Diagnose und Fix",
    title="Aus 1 Sale werden 3-5 Sales",
    left_head="So merkst du, es leckt",
    left_lines=[
        "Kundinnen kaufen einmal - dann hoerst du nichts mehr.",
        "Du hast keine Treppe (0EUR -> Mini -> Mittel -> Premium).",
        "Du machst keinen Order-Bump im Checkout.",
        "Keine Mail nach dem Kauf ausser 'Danke + Rechnung'.",
    ],
    right_head="Der Fix - die Treppe",
    right_lines=[
        "Order-Bump im Checkout (+1 kleines Produkt).",
        "Upsell direkt nach Kauf (+1 mittleres Produkt).",
        "Welcome-Sequenz fuer Kaeufer (anders als Leads).",
        "Nach 30 Tagen: Naechste Stufe der Treppe anbieten.",
    ],
    notes=(
        "DIAGNOSE + FIX VERKAUFS-LECK (5 Min):\n"
        "- Eine Kundin, die einmal gekauft hat, ist 5x wahrscheinlicher zu kaufen als ein neuer Lead.\n"
        "- Trotzdem investieren die meisten 95% in Neukunden.\n"
        "- Die Treppe muss klar sein: 0EUR -> Minikurs -> Signature -> Premium.\n"
        "- Beispiel Patricia: Bio-Check (0EUR) -> Workbook (49 CHF) -> Mentoring-Programm.\n"
        "- Hausaufgabe: Zeichnet eure Treppe auf. Welche Stufe fehlt?"
    ),
    page_num=17,
)

# === DIAGNOSE ===
slide_content(
    eyebrow="Selbst-Diagnose",
    title="Welches Leck blutet bei DIR am staerksten?",
    body_lines=[
        "Bio-Leck  -  wenn niemand reagiert obwohl Reichweite da ist.",
        "Reichweite-Leck  -  wenn deine Posts unter 500 Views bleiben.",
        "Lead-Leck  -  wenn du keine 5 neuen Email-Adressen pro Woche hast.",
        "Welcome-Leck  -  wenn deine Liste nicht kauft.",
        "Verkaufs-Leck  -  wenn deine Kundinnen nicht zurueckkommen.",
    ],
    notes=(
        "INTERAKTIVE RUNDE (8 Min):\n"
        "- Jetzt seid ihr dran. Schreibt in den Chat: WELCHES LECK ist bei euch das groesste?\n"
        "- Kurze Begruendung in 1 Satz.\n"
        "- Wir gehen die Antworten 1 fuer 1 durch und ich sage euch:\n"
        "  - Was ich als ersten Schritt empfehle\n"
        "  - In welchem Folge-Call das Leck Schwerpunkt ist."
    ),
    page_num=18,
    body_size=20,
    list_style="dot",
)

# === HAUSAUFGABE ===
slide_content(
    eyebrow="Bis zum naechsten Call",
    title="Deine Hausaufgabe - 7 Tage, 3 Schritte",
    body_lines=[
        "1. Diagnose: Welches Leck ist DEIN groesstes? (Heute entschieden.)",
        "2. EINEN Hebel anwenden - Bio neu, Hook umschreiben oder Freebie definieren.",
        "3. Live messen - 7 Tage, 1x pro Tag 5 Min in Insights schauen.",
    ],
    notes=(
        "HAUSAUFGABE (3 Min):\n"
        "- Wichtig: NUR EIN Leck angehen. Nicht alle.\n"
        "- Wer alles auf einmal aendert, kann nichts messen.\n"
        "- In Telegram-Gruppe taeglich kurz reporten - 1 Satz reicht.\n"
        "- Naechster Call: wir sprechen ueber Mail-Sequenz (Welcome-Leck).\n"
        "- Wer Fragen hat: jederzeit DM oder Telegram."
    ),
    page_num=19,
    body_size=20,
    list_style="arrow",
)

# === VORSCHAU + Closing ===
slide_content(
    eyebrow="Vorschau",
    title="Was kommt in den naechsten Calls",
    body_lines=[
        "Call 2  -  Die 5-Mail-Welcome-Sequenz, die wirklich verkauft.",
        "Call 3  -  Stories die nicht nur Likes bringen - Sales ohne Druck.",
        "Call 4  -  KI-Workflows fuer Mompreneurs - schneller, ohne Stress.",
        "Call 5  -  Skalieren ohne Burnout - das Mompreneur-Modell.",
    ],
    notes=(
        "ABSCHLUSS + VORSCHAU (3 Min):\n"
        "- Heute war die Diagnose. Jetzt geht es ans Abdichten.\n"
        "- Die naechsten 4 Calls sind alles konkrete Hebel.\n"
        "- Call 2 baut die Welcome-Sequenz LIVE - bringt eure Mail-Tools mit.\n"
        "- Termin steht im Telegram-Chat.\n"
        "- Letzte Frage in die Runde: Was nehmt ihr aus heute mit? 1 Wort reicht.\n"
        "- Danke euch. Bis naechste Woche - eure Patricia."
    ),
    page_num=20,
    body_size=20,
    list_style="dot",
)


# === Speichern ===
out = Path(__file__).resolve().parents[2] / "outputs" / "praesentationen" / "die-umsetzerinnen-call1-5-lecks.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(out)
print(f"OK -> {out}")
print(f"Folien: {len(prs.slides)}")
