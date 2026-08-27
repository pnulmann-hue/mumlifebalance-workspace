# -*- coding: utf-8 -*-
"""Baut die Präsentation für MBA Call 1 — „Es ist Arbeit. Auch wenn noch kein Geld kommt."
Haus-Stil identisch zu den Mama-CEO-Decks (Creme/Navy/Petrol/Orange, 16:9).
Sprechnotizen liegen in den Folien-Notizen — frei nachsprechen, nicht ablesen.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = r"C:\Users\pnulm\Desktop\Mein Business\outputs\produkte\mba\praesentationen"
os.makedirs(OUT, exist_ok=True)

ORANGE = RGBColor(0xDC, 0x82, 0x2E)
NAVY   = RGBColor(0x1A, 0x3A, 0x4A)
BODY   = RGBColor(0x2C, 0x3E, 0x50)
PETROL = RGBColor(0x12, 0x82, 0x8C)
MUTED  = RGBColor(0x7A, 0x8A, 0x95)
CREME  = RGBColor(0xF1, 0xEC, 0xDD)

EYEBROW = "MUM BUSINESS ACADEMY — CALL 1"
TAG     = "MBA · Umsetzerinnen-Call 1 · Dranbleiben & Konsistenz"


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = CREME


def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def run(p, text, font="Calibri", size=16, bold=False, color=BODY):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def footer(slide):
    _, tf = textbox(slide, 0.5, 5.2, 6.0, 0.3)
    run(tf.paragraphs[0], TAG, "Calibri", 9, False, MUTED)
    _, tf = textbox(slide, 7.0, 5.2, 2.5, 0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run(p, "Mum Life Balance", "Calibri", 9, False, MUTED)


def add_title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(4.2), Inches(3.5), Inches(3.5))
    c.fill.solid(); c.fill.fore_color.rgb = ORANGE; c.line.fill.background()
    _, tf = textbox(s, 0.6, 0.6, 8.8, 0.4)
    run(tf.paragraphs[0], EYEBROW, "Calibri", 11, True, ORANGE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(0.6), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    _, tf = textbox(s, 0.6, 1.35, 9.0, 2.0)
    run(tf.paragraphs[0], title, "Georgia", 34, True, NAVY)
    _, tf = textbox(s, 0.6, 3.5, 6.6, 1.0)
    run(tf.paragraphs[0], subtitle, "Georgia", 18, False, BODY)
    # Name + URL bleiben links gestapelt, damit nichts auf dem orangen Kreis liegt
    _, tf = textbox(s, 0.6, 4.62, 4.0, 0.4)
    run(tf.paragraphs[0], "Patricia Ulmann", "Calibri", 12, True, PETROL)
    _, tf = textbox(s, 0.6, 4.98, 4.0, 0.4)
    run(tf.paragraphs[0], "mumlifebalance.ch", "Calibri", 11, False, MUTED)
    return s


def add_content_slide(prs, title, bullets, mark="✓"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _, tf = textbox(s, 0.6, 0.5, 8.8, 1.0)
    run(tf.paragraphs[0], title, "Georgia", 30, True, NAVY)
    _, tf = textbox(s, 0.6, 1.55, 8.8, 3.45)
    first = True
    prev = None
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        # "~" am Anfang = Fortsetzungszeile: kein eigenes Zeichen, eingerückt
        if b.startswith("~"):
            if prev is not None:
                prev.space_after = Pt(1)
            run(p, "      ", "Calibri", 16, False, BODY)
            run(p, b[1:], "Calibri", 16, False, BODY)
        else:
            run(p, mark + "  ", "Calibri", 16, True, PETROL)
            run(p, b, "Calibri", 16, False, BODY)
        prev = p
    footer(s)
    return s


def add_statement_slide(prs, statement, source=""):
    """Grosse Aussage, zentriert — für die Kernsätze."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(0.09), Inches(2.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    _, tf = textbox(s, 1.1, 1.5, 8.0, 2.6)
    run(tf.paragraphs[0], statement, "Georgia", 30, True, NAVY)
    if source:
        _, tf = textbox(s, 1.1, 4.3, 8.0, 0.5)
        run(tf.paragraphs[0], source, "Calibri", 14, False, PETROL)
    footer(s)
    return s


def add_work_slide(prs, title, lines, hint=""):
    """Arbeits-Folie — bleibt stehen, während die Gruppe schreibt."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.5), Inches(1.75), Inches(0.42))
    badge.fill.solid(); badge.fill.fore_color.rgb = ORANGE; badge.line.fill.background()
    btf = badge.text_frame; btf.word_wrap = False
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    run(bp, "JETZT DU", "Calibri", 12, True, RGBColor(0xFF, 0xFF, 0xFF))
    _, tf = textbox(s, 0.6, 1.1, 8.8, 0.9)
    run(tf.paragraphs[0], title, "Georgia", 30, True, NAVY)
    _, tf = textbox(s, 0.6, 2.15, 8.8, 2.6)
    first = True
    for i, line in enumerate(lines, start=1):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(12)
        run(p, f"{i}.  ", "Georgia", 17, True, ORANGE)
        run(p, line, "Calibri", 17, False, BODY)
    if hint:
        _, tf = textbox(s, 0.6, 4.75, 8.8, 0.4)
        run(tf.paragraphs[0], hint, "Calibri", 13, False, PETROL)
    footer(s)
    return s


def build(filename, slides):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    for sl in slides:
        kind = sl.get("k", "c")
        if kind == "t":
            s = add_title_slide(prs, sl["t"], sl.get("sub", ""))
        elif kind == "s":
            s = add_statement_slide(prs, sl["t"], sl.get("src", ""))
        elif kind == "w":
            s = add_work_slide(prs, sl["t"], sl["b"], sl.get("hint", ""))
        else:
            s = add_content_slide(prs, sl["t"], sl["b"], sl.get("mark", "✓"))
        if sl.get("n"):
            s.notes_slide.notes_text_frame.text = sl["n"]
    prs.save(os.path.join(OUT, filename))
    print("saved", filename, len(slides), "Folien")


SLIDES = [
    {"k": "t",
     "t": "Es ist Arbeit.\nAuch wenn noch kein Geld kommt.",
     "sub": "Dranbleiben & Konsistenz \u00b7 unser erster Call",
     "n": "Herzlich willkommen zu unserem allerersten Call. Ich hab lange überlegt, womit ich starte, und ich hab mich bewusst gegen Strategie entschieden. Wir reden heute über das, was die meisten aus dem Business rauswirft, bevor es überhaupt richtig angefangen hat. Und ich fang mit etwas an, das ich lange niemandem erzählt hab."},

    {"k": "c", "t": "Das machen wir heute",
     "b": ["Kurze Runde \u2014 wer bist du, seit wann machst du das",
           "Warum du dich am falschen Ergebnis misst",
           "Deine Bestandsaufnahme: was du l\u00e4ngst gebaut hast",
           "Warum vier Wochen die Mindestdosis sind",
           "Deine eine Sache f\u00fcr die n\u00e4chsten vier Wochen"],
     "n": "Kurz, was heute passiert, damit du weisst, worauf es hinausläuft. Wir machen eine Runde, dann rede ich einen Moment, danach arbeitest du selber, dann rede ich nochmal, und dann arbeitest du nochmal. Am Ende gehst du mit einer ganz konkreten Sache raus. Das ist mir wichtig, ich will keine Calls machen, nach denen man ein gutes Gefühl hat und am nächsten Tag nichts anders ist. Und eins sag ich gleich vorneweg, damit es keine falsche Erwartung gibt: es geht heute um Dranbleiben, aber es geht nicht um Vollgas. Wenn bei dir gerade wenig oder gar nichts geht, dann bist du hier genau richtig, denn genau dafür bauen wir das nachher passend."},

    {"k": "w", "t": "Kurze Runde",
     "b": ["Dein Name und was du verkaufst",
           "Seit wann machst du das eigentlich schon?"],
     "hint": "Die zweite Frage ist mir heute die wichtigste \u2014 du wirst gleich merken warum.",
     "n": "Fangen wir mit einer Runde an. Sag mir deinen Namen, was du verkaufst, und dann die Frage, die mir heute die wichtigste ist: seit wann machst du das eigentlich schon. Ganz kurz, ein Satz reicht, wir gehen nachher tief genug rein. Ich schreib mit."},

    {"k": "s",
     "t": "\u201eIch hab jahrelang gedacht,\ndas hier sei gar keine\nrichtige Arbeit.\u201c",
     "src": "Weil kein Geld reingekommen ist.",
     "n": "Jetzt mal ganz ehrlich. Ich hab in den ersten Jahren gedacht, das was ich hier mache sei gar keine richtige Arbeit. Weil kein Geld reingekommen ist. Ich hab am Abend am Küchentisch gesessen, stundenlang was aufgebaut, und dann kam dieses Gefühl: ja was machst du hier eigentlich, das bringt ja nichts. Ich hab mich selber nicht ernst genommen, obwohl ich jeden einzelnen Tag gearbeitet hab. Und ich sag dir das, weil ich ziemlich sicher bin, dass mindestens eine von euch genau da gerade sitzt."},

    {"k": "c", "t": "Der Denkfehler", "mark": "–",
     "b": ["Du schaust auf ein einziges Ergebnis: was reinkommt",
           "Und das ist genau das, was am l\u00e4ngsten braucht",
           "Alles andere, was du aufbaust, z\u00e4hlst du gar nicht",
           "Also kommst du logischerweise zum Schluss: bringt nichts"],
     "n": "Der Denkfehler dahinter ist eigentlich ganz einfach. Du schaust auf ein einziges Ergebnis, nämlich auf das, was reinkommt. Und das ist ausgerechnet das, was am längsten braucht. Alles andere, was du in der Zwischenzeit aufbaust, zählst du gar nicht mit, weil es keine Zahl auf dem Konto ist. Und dann kommst du logischerweise zum Schluss, dass das nichts bringt. Dabei stimmt die Rechnung einfach nicht."},

    {"k": "c", "t": "Du hast vier Konten, nicht eins",
     "b": ["Sichtbarkeit \u2014 wer dich \u00fcberhaupt kennt",
           "Vertrauen \u2014 wer dir zuh\u00f6rt und antwortet",
           "Aufgebautes \u2014 was bleibt, auch wenn du zwei Wochen krank bist",
           "Geld \u2014 was reinkommt"],
     "n": "Ich zeig dir, wie ich das heute sehe. Du hast nämlich nicht ein Konto, sondern vier. Das erste ist Sichtbarkeit, also wer dich überhaupt kennt. Das zweite ist Vertrauen, also wer dir zuhört und antwortet. Das dritte ist alles, was du aufgebaut hast und was bleibt, auch wenn du mal zwei Wochen flach liegst. Und das vierte ist Geld. Vier Konten, und du schaust nur auf eins davon."},

    {"k": "c", "t": "Woran du sie merkst",
     "b": ["Sichtbarkeit: Aufrufe, neue Followerinnen, gespeicherte Beitr\u00e4ge",
           "Vertrauen: Antworten auf Stories, DMs, Menschen auf deiner Liste",
           "Aufgebautes: Profil, Freebie, Liste, Vorlagen, dein Wissen",
           "\u201eIch verfolge dich schon l\u00e4nger\u201c \u2014 das ist eine Einzahlung"],
     "n": "Und damit das nicht schwammig bleibt, hier ganz konkret, woran du die merkst. Sichtbarkeit sind Aufrufe, neue Followerinnen, gespeicherte Beiträge. Vertrauen sind Antworten auf deine Stories, DMs, Leute auf deiner Liste. Aufgebautes ist dein Profil, dein Freebie, deine Vorlagen, dein Wissen. Und wenn dir jemand schreibt, ich verfolge dich schon länger, dann ist das eine Einzahlung. Die verbuchst du nur nie."},

    {"k": "s",
     "t": "Geld ist das Konto,\ndas sich zuletzt f\u00fcllt.",
     "src": "Nicht das, an dem du dich zuerst messen darfst.",
     "n": "Und darum ist mein Satz der hier: Geld ist das Konto, das sich zuletzt füllt. Es wird nämlich aus den ersten drei gespeist. Niemand kauft bei jemandem, den sie nicht kennt, dem sie nicht vertraut und der nichts hat, was sie kaufen könnte. Wenn du also nur auf Konto vier schaust, schaust du aufs Ergebnis von einer Arbeit, die du gerade erst angefangen hast."},

    {"k": "c", "t": "Wie lange es wirklich dauert",
     "b": ["Meine Erfahrung: online mindestens ein Jahr",
           "Und zwar unter der Bedingung, dass du konstant dranbleibst",
           "Bei mir hat es l\u00e4nger gedauert \u2014 selbst\u00e4ndig seit 2023,",
           "~regelm\u00e4ssiges Einkommen ab Fr\u00fchling 2025"],
     "n": "Jetzt zur ehrlichen Zahl, und das ist keine Statistik, das ist meine Erfahrung. Online brauchst du normalerweise mindestens ein Jahr, bis regelmässig Geld kommt. Und zwar unter der Bedingung, dass du konstant dranbleibst. Bei mir hat es sogar länger gedauert, ich bin seit 2023 selbständig und regelmässig kam es ab Frühling 2025. Ich sag dir das nicht, um dich zu entmutigen, sondern damit du in Woche sechs nicht denkst, du machst was falsch."},

    {"k": "w", "t": "Deine Bestandsaufnahme",
     "b": ["Was hast du in den letzten zw\u00f6lf Monaten aufgebaut,\n     das vorher nicht da war?",
           "Wer kennt dich heute, der dich vor einem Jahr nicht kannte?",
           "Was kannst du heute, was du vor einem Jahr nicht konntest?"],
     "hint": "Alles z\u00e4hlt. Auch ein Kurs, den du durchgearbeitet hast.",
     "n": "So, jetzt bist du dran. Du hast sieben, acht Minuten und schreibst zu diesen drei Fragen auf, was dir einfällt. Und alles zählt: dein Profil, deine Beiträge, dein Freebie, deine Liste, ein Kurs, den du durchgearbeitet hast, ein Thema, das klarer geworden ist. Danach lesen wir reihum je zwei Sachen vor, und ich sag dir bei jeder, auf welches Konto das eingezahlt hat."},

    {"k": "s",
     "t": "Und dann gibt es Phasen,\nin denen gar nichts geht.",
     "src": "Auch die gehören dazu. Auch die sind Teil vom Aufbau.",
     "n": "Und bevor wir zum zweiten Teil kommen, muss ich noch etwas sagen, weil es sonst schief ankommt. Es gibt Phasen, in denen gar nichts geht. Wo du keine Kraft hast, wo privat etwas läuft, wo dein Business einfach steht, wochenlang. Und ich will, dass ihr das von mir hört und nicht von irgendeinem Coach im Internet: das gehört dazu. Das ist kein Rückschritt und es ist auch nicht dein Versagen. Ich hatte solche Phasen und ich werde wieder welche haben."},

    {"k": "c", "t": "Was in so einer Phase gilt",
     "b": ["Eine Pause ist keine verlorene Zeit — dein Aufgebautes bleibt",
           "Pause ist nicht aufhören — aufhören heisst stilles Wegbleiben",
           "Sag es laut, dann muss niemand raten — ein Satz bei uns reicht",
           "Und wenn du wieder kannst, fängst du nicht bei null an"],
     "n": "Vier Sachen, die in so einer Phase gelten. Erstens: eine Pause ist keine verlorene Zeit, weil dein aufgebautes Konto bleibt. Dein Profil, dein Wissen, deine Leute, das verschwindet nicht in sechs Wochen. Zweitens, und das ist mir das wichtigste: Pause ist nicht aufhören. Aufhören heisst wegbleiben, ohne dass es jemand merkt. Drittens: sag es laut, ein Satz bei uns im Fragen-Bereich reicht, dann muss niemand raten und niemand denkt, du hättest es aufgegeben. Und viertens: wenn du wieder kannst, fängst du nicht bei null an. Du fängst da an, wo du aufgehört hast, und das ist ein grosser Unterschied."},

    {"k": "s",
     "t": "Vier Wochen sind\nkeine Challenge.",
     "src": "Vier Wochen sind die Mindestdosis, bevor \u00fcberhaupt etwas messbar wird.",
     "n": "Und damit sind wir beim zweiten Teil, denn Anerkennung allein bringt dich nicht weiter. Jetzt kommt die Konsistenz. Und mein Satz dazu: vier Wochen sind keine Challenge, vier Wochen sind die Mindestdosis. Vorher kannst du gar nicht beurteilen, ob etwas funktioniert."},

    {"k": "c", "t": "Was ich am h\u00e4ufigsten sehe", "mark": "\u2013",
     "b": ["Jemand zieht etwas richtig durch \u2014 und es funktioniert sofort",
           "Die Reichweite geht rauf, es kommt Bewegung rein",
           "Dann kommt eine Woche mit krankem Kind",
           "Und danach kommt es nie wieder zur\u00fcck \u2014 alles f\u00e4llt zusammen"],
     "n": "Ich erzähl dir, was ich am häufigsten sehe. Und ich hab da ein Beispiel aus meiner Begleitung: eine Frau hat einen Monat lang jeden Tag gepostet. Die Reichweite ist deutlich hochgegangen, es kam richtig Bewegung rein, es hat funktioniert. Das ist der Beweis, dass Konstanz wirkt, und zwar schneller, als die meisten glauben. Und dann hat sie aufgehört, und ein paar Wochen später war alles wieder da, wo es vorher war."},

    {"k": "s",
     "t": "Es liegt nie an dir als Person.\nMeistens ist einfach\ndas System zu gross gebaut.",
     "src": "Jeden Tag ein Beitrag neben Kindern ist kein Rhythmus. Das ist ein Sprint.",
     "n": "Und jetzt ganz wichtig, damit das nicht falsch ankommt: das ist kein Versagen. Es liegt nie an dir als Person, meistens ist einfach das System zu gross gebaut. Jeden Tag ein Beitrag neben Kindern ist kein Rhythmus, das ist ein Sprint. Und Sprints hält niemand über Monate durch. Nicht die Intensität entscheidet, sondern ob du es noch machst, wenn es langweilig geworden ist."},

    {"k": "c", "t": "Die Regel, die alles \u00e4ndert",
     "b": ["W\u00e4hl EINE Sache. Nicht drei.",
           "Sie muss so klein sein, dass du sie auch an einem",
           "~beschissenen Tag noch schaffst",
           "Und sie kriegt eine Notfall-Version \u2014 gleichwertig, kein R\u00fcckfall"],
     "n": "Darum machen wir das heute anders. Du suchst dir eine Sache. Nicht drei, eine. Und die muss so klein sein, dass du sie auch an dem Tag noch schaffst, an dem alles schiefgeht. Und dann kriegt sie eine Notfall-Version. Und jetzt hör mir gut zu: die Notfall-Version ist kein Rutsch nach unten und keine schlechtere Variante. Sie ist deine Sache, nur kleiner. Wer die Notfall-Version macht, hat die Woche geschafft, Punkt."},

    {"k": "c", "t": "So sieht das aus",
     "b": ["Jeden Tag eine Story  \u2192  eine einzige Story mit dem Handy im Auto",
           "Drei Beitr\u00e4ge pro Woche  \u2192  einer, daf\u00fcr fix am selben Tag",
           "Jeden Tag zwei echte DMs  \u2192  eine Sprachnachricht an eine Person",
           "Sonntags die Woche planen  \u2192  zehn Minuten, drei Zeilen, fertig"],
     "n": "Ein paar Beispiele, damit du ein Gefühl kriegst. Links steht deine Sache, rechts die Notfall-Version. Jeden Tag eine Story wird zu einer einzigen Story mit dem Handy im Auto. Drei Beiträge die Woche werden zu einem, dafür fix am selben Tag. Und die Sonntagsplanung wird zu zehn Minuten, drei Zeilen, fertig. Du merkst: die Notfall-Version darf richtig klein sein. Sie muss nur stattfinden."},

    {"k": "w", "t": "Deine eine Sache",
     "b": ["Meine eine Sache f\u00fcr die n\u00e4chsten vier Wochen: \u2026",
           "Meine Notfall-Version f\u00fcr schlechte Tage: \u2026"],
     "hint": "Zwei Zeilen \u2014 danach sagst du sie laut in die Runde.",
     "n": "Jetzt schreibst du zwei Zeilen auf: deine eine Sache für die nächsten vier Wochen und deine Notfall-Version für schlechte Tage. Danach sagt sie jede laut, und ich sag dir ehrlich, wenn ich sie zu gross finde. Frag dich beim Aufschreiben schon selber: und am Tag, an dem beide Kinder krank sind, schaff ich das dann noch?"},

    {"k": "c", "t": "Deine Aufgabe f\u00fcr vier Wochen",
     "b": ["Vier Wochen deine eine Sache durchziehen",
           "Jede Woche drei Zeilen bei uns im Fragen-Bereich",
           "Meine Sache \u00b7 Durchgezogen: ja / mit Notfall-Version / nein",
           "Was passiert ist — und da geht es NICHT ums Geld",
           "Geht gerade gar nichts? Dann SIND die drei Zeilen deine eine Sache"],
     "n": "Deine Aufgabe ist einfach und trotzdem die schwerste, die ich dir dieses Jahr gebe. Du ziehst deine eine Sache vier Wochen durch, mit Notfall-Version wenn es sein muss. Und du schreibst jede Woche drei Zeilen bei uns in den Fragen-Bereich. Was du geschafft hast, ob mit Notfall-Version, und was passiert ist. Und da geht es ausdrücklich nicht ums Geld, sondern um alles andere: wer sich gemeldet hat, was du fertig gebaut hast, was leichter ging als letzte Woche."},

    {"k": "c", "t": "Und wenn eine Woche nicht klappt?",
     "b": ["Dann schreibst du genau das rein",
           "Wir schauen gemeinsam, ob deine Sache zu gross gebaut war",
           "Dranbleiben heisst nicht l\u00fcckenlos \u2014 es heisst zur\u00fcckkommen",
           "Pause ist nicht aufhören — aufhören heisst stilles Wegbleiben"],
     "n": "Und wenn eine Woche nicht klappt, dann schreibst du genau das rein. Nein ist eine völlig gültige Antwort, und dann schauen wir gemeinsam, ob deine Sache zu gross gebaut war, und machen sie kleiner. Dranbleiben heisst nämlich nicht lückenlos, es heisst zurückkommen. Und wenn es länger dauert als eine Woche, dann ist das eine Pause, und Pause ist nicht aufhören. Aufhören heisst wegbleiben, ohne dass es jemand merkt — und genau das musst du hier nie machen."},

    {"k": "c", "t": "Wie es weitergeht",
     "b": ["Zwei Calls im Monat \u2014 ein Thema mit euren Fragen, einer nur f\u00fcr Fragen",
           "Alles wird aufgezeichnet und liegt in der Academy",
           "Zwischendurch: Fragen-Bereich, ich bin Di & Do drin",
           "N\u00e4chster Themen-Call: deine vier Wochen \u2014 und dein Thema"],
     "n": "Zum Schluss noch kurz, wie es weitergeht. Wir haben zwei Calls im Monat, einen zu einem Thema mit euren Fragen und einen nur für eure Fragen. Alles wird aufgezeichnet und liegt danach in der Academy, du verpasst also nichts, wenn ein Kind Fieber hat. Zwischendurch läuft der Austausch im Fragen-Bereich und ich bin dienstags und donnerstags drin. Und beim nächsten Themen-Call schauen wir zwei Sachen an: was aus deinen vier Wochen geworden ist, und dann nehmen wir uns dein Thema vor. Ich freu mich auf euch."},
]

build("call-01-dranbleiben.pptx", SLIDES)
print("ALL DONE")
