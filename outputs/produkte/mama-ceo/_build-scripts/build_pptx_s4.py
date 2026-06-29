# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r"C:\Users\pnulm\Desktop\Mein Business\outputs\produkte\mama-ceo\03-praesentationen\saeule-4"
os.makedirs(OUT, exist_ok=True)

ORANGE = RGBColor(0xDC, 0x82, 0x2E)
NAVY   = RGBColor(0x1A, 0x3A, 0x4A)
BODY   = RGBColor(0x2C, 0x3E, 0x50)
PETROL = RGBColor(0x12, 0x82, 0x8C)
MUTED  = RGBColor(0x7A, 0x8A, 0x95)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREME  = RGBColor(0xF1, 0xEC, 0xDD)

def set_bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = CREME

EYEBROW = "SÄULE 4 — DU DELEGIERST DEN ADMINKRAM"

def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf

def run(p, text, font="Calibri", size=16, bold=False, color=BODY):
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color
    return r

def add_title_slide(prs, title, subtitle, eyebrow=EYEBROW):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(s)
    # deko-kreis unten rechts
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(4.2), Inches(3.5), Inches(3.5))
    c.fill.solid(); c.fill.fore_color.rgb = ORANGE; c.line.fill.background()
    # eyebrow
    _, tf = textbox(s, 0.6, 0.6, 8.8, 0.4)
    p = tf.paragraphs[0]; run(p, eyebrow, "Calibri", 11, True, ORANGE)
    # accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(0.6), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    # title
    _, tf = textbox(s, 0.6, 1.35, 8.6, 1.7)
    p = tf.paragraphs[0]; run(p, title, "Georgia", 44, True, NAVY)
    # subtitle
    _, tf = textbox(s, 0.6, 3.15, 8.4, 1.4)
    p = tf.paragraphs[0]; run(p, subtitle, "Georgia", 18, False, BODY)
    # footer
    _, tf = textbox(s, 0.6, 5.0, 4.0, 0.4)
    p = tf.paragraphs[0]; run(p, "Patricia Ulmann", "Calibri", 12, True, PETROL)
    _, tf = textbox(s, 5.4, 5.0, 4.0, 0.4)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; run(p, "mumlifebalance.ch", "Calibri", 11, False, MUTED)
    return s

def add_content_slide(prs, title, bullets, tag):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _, tf = textbox(s, 0.6, 0.5, 8.8, 1.0)
    p = tf.paragraphs[0]; run(p, title, "Georgia", 30, True, NAVY)
    # body
    _, tf = textbox(s, 0.6, 1.55, 8.8, 3.45)
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        run(p, "✓  ", "Calibri", 16, True, PETROL)
        run(p, b, "Calibri", 16, False, BODY)
    # footer
    _, tf = textbox(s, 0.5, 5.2, 6.0, 0.3)
    p = tf.paragraphs[0]; run(p, tag, "Calibri", 9, False, MUTED)
    _, tf = textbox(s, 7.0, 5.2, 2.5, 0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; run(p, "Mum Life Balance", "Calibri", 9, False, MUTED)
    return s

def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def build(filename, lesson_tag, slides):
    prs = Presentation()
    prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    for i, sl in enumerate(slides):
        if i == 0:
            s = add_title_slide(prs, sl["t"], sl.get("sub", ""))
        else:
            s = add_content_slide(prs, sl["t"], sl["b"], lesson_tag)
        if sl.get("n"):
            set_notes(s, sl["n"])
    path = os.path.join(OUT, filename)
    prs.save(path)
    print("saved", filename, len(slides), "Folien")

# ============ LESSON DATA ============
L41 = [
 {"t":"KI-Mythos vs. Realität","sub":"Input ist alles","n":"Willkommen in Säule 4, dem Teil, auf den ich mich am meisten freue, weil hier deine KI-Mitarbeiter entstehen. Aber bevor wir irgendwas bauen, müssen wir über die eine Sache reden, die fast jede Mama falsch macht, wenn sie mit KI startet."},
 {"t":"„KI bringt mir nichts\"","b":["Aufgemacht, gefragt, unbrauchbar — wieder zu","Das kennen viele","Es liegt fast nie an der KI"],"n":"Hand aufs Herz, vielleicht hast du es selbst schon gedacht: ich hab ChatGPT mal aufgemacht, irgendwas gefragt, und das Ergebnis war so allgemein und unbrauchbar, dass ich es wieder zugemacht hab. Das kenn ich, und ich verspreche dir, es liegt fast nie an der KI."},
 {"t":"Der Mythos: KI als Zauberkasten","b":["❌ Knopf drücken → fertiges Business","So funktioniert das nicht","Die Wahrheit: DU bleibst am Steuer"],"n":"Der grösste Irrtum ist, dass KI ein Zauberkasten ist, wo du auf einen Knopf drückst und dein halbes Business kommt fertig raus. So funktioniert das nicht, und wer das erwartet, ist enttäuscht. Die Wahrheit ist viel beruhigender, weil sie bedeutet, dass DU am Steuer bleibst."},
 {"t":"KI ist eine Praktikantin","b":["Schnell, fleissig, nie müde","Weiss am ersten Tag nichts über dich","Nur gut mit Einarbeitung"],"n":"Stell dir KI lieber als Praktikantin vor. Sie ist superschnell, sie wird nie müde, und sie beschwert sich nie. Aber sie weiss am ersten Tag absolut nichts über dich, dein Business und deine Kundinnen. Und genau wie eine echte Praktikantin liefert sie nur dann was Brauchbares, wenn du sie vernünftig einarbeitest."},
 {"t":"Müll rein, Müll raus","b":["Vage Frage → vage Antwort","Klarer Auftrag + wer du bist → brauchbar","Dein Ergebnis = so gut wie dein Auftrag"],"n":"Daraus folgt die wichtigste KI-Regel überhaupt: Müll rein, Müll raus. Wenn du vage fragst, kriegst du Vages zurück, und wenn du genau sagst was du brauchst und wer du bist, kriegst du was Brauchbares. Dein Ergebnis ist immer nur so gut wie dein Auftrag, und das ist eine gute Nachricht, weil du den Auftrag in der Hand hast."},
 {"t":"Warum meine Helfer funktionieren","b":["Kein Tech-Studium — Mama von vier","Alles selbst beigebracht","Es liegt an guten Aufträgen — lernbar"],"n":"Ich sag dir mal ganz ehrlich, warum bei mir so viele KI-Helfer laufen: nicht weil ich programmieren kann oder ein Tech-Genie bin, ich bin Mama von vier Kindern und hab mir das alles selbst beigebracht. Es funktioniert, weil ich gelernt hab, gute Aufträge zu geben. Das kannst du genauso lernen, und genau das machen wir in dieser Säule."},
 {"t":"3 Dinge für jeden guten Auftrag","b":["1. Kontext — wer bin ich","2. Klare Aufgabe — was genau","3. Beispiel — wie soll's aussehen"],"n":"Ein guter KI-Auftrag hat immer drei Zutaten. Erstens Kontext, also wer bist du und für wen arbeitest du. Zweitens eine klare Aufgabe, nicht „mach mir Content\", sondern „schreib mir drei Story-Ideen zum Thema X für Mamas im Network\". Und drittens ein Beispiel, damit sie deinen Stil trifft. Diese drei Dinge bauen wir in den nächsten Lektionen Stück für Stück auf."},
 {"t":"Der Fahrplan für Säule 4","b":["🌅 Cockpit-Bot — Business-Morgenbriefing","🏠 Haushalts-Helfer — Mental Load zu Hause","Zwillinge: gleiche Mechanik"],"n":"Und das ist unser Fahrplan: wir bauen zwei KI-Mitarbeiter, die im Grunde Zwillinge sind. Den Cockpit-Bot, der morgens in dein Notion-Business schaut, und den Haushalts-Helfer, der genauso deine Haushalts-Liste im Kopf behält. Einer fürs Business, einer fürs Zuhause, weil der Adminkram an beiden Orten an dir zerrt."},
 {"t":"Arbeitsblatt + nächste Lektion","b":["📋 Notier deinen letzten frustrierenden KI-Versuch","Nächste Lektion: 4.2 — womit du arbeitest"],"n":"Im Arbeitsblatt schreibst du kurz auf, wo dich KI bisher genervt hat, denn am Ende der Säule schaust du nochmal drauf und siehst den Unterschied. In der nächsten Lektion klären wir kurz, womit du überhaupt arbeitest, damit du nicht in der Auswahl hängenbleibst. Bis gleich."},
]

L42 = [
 {"t":"Womit du arbeitest","sub":"Die 3 Stufen — Cowork ist dein Freund","n":"In dieser Lektion klären wir, womit du überhaupt arbeitest, und ich halt das ganz einfach. Ich zeig dir drei Stufen, und das Schöne ist: auf jeder nimmst du was mit, auch wenn du nie irgendwas installierst. Du musst nur wissen, auf welcher du starten willst."},
 {"t":"Der nackte Chat reicht nicht","b":["Gratis-Chat vergisst dich jedes Mal","Zu statisch fürs Business","Du willst was, das dein Wissen behält"],"n":"Das Erste, was du wissen musst: das normale Gratis-Chatfenster, wo du was reintippst und morgen ist alles weg, das reicht fürs Business auf Dauer nicht. Es ist zu statisch und vergisst dich jedes Mal. Du willst was, das dein Wissen behält und in dein Notion schaut, und dafür zeig ich dir die drei Stufen."},
 {"t":"Die 3 Stufen auf einen Blick","b":["🟢 Stufe 0 — nur Notion (kein Bot)","🟡 Stufe 1 — Claude Cowork (kein Code)","🔵 Stufe 2 — Claude Code (automatisch)"],"n":"Stufe null ist nur dein Notion, ganz ohne Bot. Stufe eins ist Claude Cowork, ein Bot ohne Code. Stufe zwei ist Claude Code, der automatische Profi-Bot. Du fängst da an, wo du dich wohlfühlst, und kannst jederzeit eine Stufe höher. Niemand muss bis Stufe zwei, ganz im Gegenteil."},
 {"t":"Stufe 0: Notion allein ist schon ein Win","b":["Liste + fertige Ansichten öffnen","Du siehst Termine + was am Tag dran ist","Kein Bot, kein Abo — Kopf ist frei"],"n":"Stufe null heisst: du hast deine Liste in Notion, mit ein paar fertigen Ansichten, und machst sie einfach auf. Du siehst, welche Termine kommen und was an welchem Tag dran ist. Kein Bot, kein Abo, gar nichts. Allein dass das Zeug aus deinem Kopf raus und an einem Ort ist, ist schon die halbe Miete."},
 {"t":"Stufe 1: Claude Cowork — dein Bot ohne Code","b":["Desktop-App mit Knöpfen, kein Terminal","Notion verbinden + Vorlage einfügen","Fragen: „was ist heute dran?\""],"n":"Stufe eins ist mein Liebling für die meisten von euch: Claude Cowork. Das ist eine App auf deinem Computer mit Knöpfen, kein Terminal, kein Programmieren. Du verbindest einmal dein Notion, fügst meine fertige Bot-Vorlage ein, und ab dann fragst du „was ist heute dran\" und kriegst deine Antwort aus deinem echten Notion. So einfach ist das."},
 {"t":"Was ist ein „Connector\"?","b":["Der Stecker zwischen Claude und Notion","Einmal einstecken (Notion-Login)","Du gibst frei, was er sehen darf"],"n":"Kurz zum Wort Connector, weil's wichtig ist: stell dir das wie einen Stecker vor. Notion liefert den Stecker, und in Cowork klickst du ihn einmal ein, loggst dich bei Notion ein und sagst, welche Seiten Claude sehen darf. Ab dann reden die zwei miteinander. Genau das kann der nackte Gratis-Chat nicht, deshalb brauchst du Cowork oder Claude Code dafür."},
 {"t":"Stufe 2: Claude Code (Ausblick)","b":["Bot läuft von selbst, aufs Handy","Auch wenn der Laptop zu ist","Profi-Stufe, kein Muss"],"n":"Stufe zwei ist Claude Code, damit baust du einen Bot, der von selbst läuft und dir morgens aufs Handy schreibt, auch wenn dein Laptop zu ist. So läuft meiner. Das ist die Profi-Stufe, braucht etwas mehr und ist ausdrücklich kein Muss. Ich zeig's dir in 4.4 als Ausblick, aber die allermeisten bleiben glücklich auf Stufe eins."},
 {"t":"Was kostet das?","b":["Stufe 0: gratis","Stufe 1 Cowork: ~20-23 CHF/Monat (Claude Pro)","Stufe 2: + Hosting ~5 CHF/Monat"],"n":"Zum Geld, ganz ehrlich: Stufe null ist gratis. Für Cowork brauchst du ein Claude-Pro-Abo, etwa zwanzig bis dreiundzwanzig Franken im Monat, das ist weniger als ein Mittagessen auswärts und nimmt dir Stunden ab. Nur die Profi-Stufe kostet später noch ein bisschen Hosting obendrauf, aber dazu kommen wir in 4.4."},
 {"t":"Eins nach dem anderen","b":["Nicht 5 Sachen gleichzeitig","Eine Stufe, zwei Bots, dann weiterschauen","Sonst verzettelst du dich"],"n":"Ein Rat aus Erfahrung: fang nicht mit fünf Sachen gleichzeitig an. Wähl eine Stufe, bau damit deine zwei Helfer aus dieser Säule, gewöhn dich dran, und erst wenn das sitzt, schaust du dir die nächste an. Sonst verzettelst du dich, und das ist genau das Hamsterrad, aus dem wir dich rausholen."},
 {"t":"Arbeitsblatt + nächste Lektion","b":["📋 Wähle deine Stufe + leg den Zugang an","Nächste Lektion: 4.3 — dein Business-Brief"],"n":"Im Arbeitsblatt entscheidest du jetzt, auf welcher Stufe du startest, und legst, falls nötig, deinen Zugang an, damit du in 4.4 sofort bauen kannst. In der nächsten Lektion schreiben wir deinen Business-Brief, das Herzstück, das jeder deiner Helfer braucht. Bis gleich."},
]

L43 = [
 {"t":"Dein Business-Brief","sub":"Das Wissen, das jeder Bot braucht","n":"Jetzt bauen wir das Herzstück deiner ganzen KI-Arbeit, und das ist nicht ein Bot, sondern ein kurzes Dokument, das alle deine Bots klüger macht. Ich nenne es den Business-Brief, und wenn der einmal steht, wird alles andere leicht."},
 {"t":"Warum jeder Bot dasselbe Wissen braucht","b":["Nicht jeder Helfer einzeln erklären","Einmal aufschreiben → jede KI kriegt es","Du briefst einmal, nicht zehnmal"],"n":"Stell dir vor, du holst dir zwei Praktikantinnen ins Haus. Du würdest doch nicht jeder einzeln und jeden Tag neu erklären, wer du bist und was du machst. Genau deshalb schreibst du dein Wissen einmal auf, und jede neue KI kriegt dasselbe Dokument. Du briefst einmal, nicht zehnmal."},
 {"t":"Ohne Brief: Einheitsbrei","b":["Ohne Brief klingt jede KI gleich","Mit Brief klingt sie wie DU","Unterschied: nutzlos vs. beste Mitarbeiterin"],"n":"Ohne diesen Brief klingt jede KI gleich, nämlich nach allgemeinem Internet-Einheitsbrei, und das spüren deine Leserinnen sofort. Mit dem Brief klingt sie wie du, mit deiner Sprache, deinen Themen, deiner Haltung. Das ist der Unterschied zwischen nutzlos und meine beste Mitarbeiterin."},
 {"t":"Die 6 Bausteine","b":["1 Wer bin ich · 2 Wer ist meine Kundin","3 Was biete ich an · 4 Meine Themen","5 Meine Stimme · 6 Was der Bot NIE tut"],"n":"Dein Business-Brief hat sechs Bausteine. Erstens: wer bin ich. Zweitens: wer ist meine Kundin. Drittens: was biete ich an, also deine Produkte. Viertens: über welche Themen rede ich. Fünftens: wie klinge ich, deine Stimme. Und sechstens, ganz wichtig: was soll der Bot NICHT tun, deine Tabus. Diese sechs füllst du gleich im Arbeitsblatt aus."},
 {"t":"Baustein 1-3: dein Fundament","b":["Wer bin ich (Mama, Network, Thema)","Kundin — konkret, nicht „alle\"","Angebot: Gratis / Mini / Gross (aus Säule 3)"],"n":"Die ersten drei Bausteine kennst du eigentlich schon aus den letzten Säulen. Wer du bist, deine Kundin so konkret wie möglich, nicht alle Frauen, und dein Angebot, das du in Säule 3 als Gratis-, Mini- und grosses Produkt schon sortiert hast. Du schreibst hier nur zusammen, was du längst weisst."},
 {"t":"Baustein 4-6: dein Charakter","b":["Themen — Überthemen aus der Jahres-Strategie","Stimme — locker/seriös, du/Sie","Tabus — was nie vorkommt"],"n":"Die zweiten drei geben dem Bot deinen Charakter. Deine Überthemen aus der Jahres-Strategie, deine Stimme, also redest du eher locker oder seriös, per du oder per Sie, und deine Tabus, zum Beispiel keine übertriebenen Versprechen, keine Fremdwörter, was auch immer dir wichtig ist. Damit hat die KI Leitplanken."},
 {"t":"Stimme einfangen: gib Beispiele","b":["Häng 2-3 echte Texte an","Caption, Nachricht — irgendwas nach dir","Die KI ahmt deinen Ton nach"],"n":"Der beste Trick für deine Stimme ist erstaunlich einfach: häng zwei, drei deiner echten Texte an den Brief, eine Caption, eine Nachricht, irgendwas das nach dir klingt. Die KI liest das und ahmt deinen Ton nach. Du musst deine Stimme nicht beschreiben können, du musst sie nur zeigen."},
 {"t":"Der einfachere Weg: lass dich interviewen","b":["Sag: „schreib ein Buch über mich\"","Die KI fragt so lange, bis sie alles hat","Fertiger Brief am Schluss — Prompt liegt als Bonus bei"],"n":"Und jetzt mein Lieblings-Trick, den ich selbst die ganze Zeit benutze: du musst dieses Blatt gar nicht alleine ausfüllen. Sag der KI einfach, sie soll tun, als würde sie ein Buch über dich und dein Business schreiben, und sie soll dir so lange und so vertieft Fragen stellen, bis sie wirklich alle Infos von dir hat. Dann sitzt du nicht vor einem leeren Blatt, sondern beantwortest einfach eine Frage nach der anderen, und am Schluss spuckt sie dir deinen fertigen Business-Brief aus. Genau diesen Prompt hab ich dir als Bonus dazugelegt, du musst ihn nur reinkopieren."},
 {"t":"Mein Business-Brief als Beispiel","b":["Nichts Hochgestochenes","Normale Worte: wer, für wen, wie, was nie","Grob und ehrlich reicht"],"n":"Ich zeig dir jetzt meinen eigenen Business-Brief, damit du ein Gefühl kriegst. Du siehst, das ist nichts Hochgestochenes, sondern in normalen Worten: wer ich bin, für wen ich da bin, wie ich rede und was bei mir nie vorkommt. Genau so grob und ehrlich darf deiner auch sein."},
 {"t":"Arbeitsblatt + nächste Lektion","b":["📋 Schreib deinen Business-Brief (6 Bausteine)","Nächste Lektion: 4.4 — Cockpit-Bot bauen"],"n":"Im Arbeitsblatt schreibst du jetzt deinen Business-Brief, Baustein für Baustein, und du legst ihn als Notiz ab, wo du ihn immer findest. Den brauchst du in der nächsten Lektion sofort, denn da bauen wir deinen ersten echten Bot, den Cockpit-Bot. Bis gleich."},
]

L44 = [
 {"t":"Dein Cockpit-Bot","sub":"Dein persönliches Morgenbriefing","n":"Jetzt wird's konkret, denn wir bauen deinen ersten echten KI-Mitarbeiter. Ich nenne ihn den Cockpit-Bot, weil er dir morgens das Cockpit zeigt, also alles was heute wichtig ist, auf einen Blick."},
 {"t":"Was er für dich tut","b":["Morgens fragen: „was ist heute dran?\"","→ Tagesfokus + 3 Hauptaufgaben + Wochenblick","Kein Wühlen in Notion mehr"],"n":"Was dieser Bot macht, ist einfach erklärt: du fragst ihn morgens „was ist heute dran\", und er gibt dir deinen Tagesfokus, deine drei wichtigsten Aufgaben und einen kurzen Überblick über die Woche. Du musst nicht mehr selbst in Notion wühlen und dich verlieren, er macht das für dich."},
 {"t":"Warum wir Säule 3 gebraucht haben","b":["Bot liest dein Notion-Business-Brain","Wochenplan, Tagesplaner, Aufgaben","Ohne Brain kein Briefing"],"n":"Jetzt zahlt sich Säule 3 aus. Dein Cockpit-Bot kann nur deshalb wissen was heute dran ist, weil du dein Notion-Business-Brain aufgebaut hast, mit Wochenplan, Tagesplaner und Aufgaben. Der Bot ist die Stimme, die dir vorliest, was in deinem Brain steht. Ohne das Brain kein Briefing."},
 {"t":"3 Stufen — wir bauen Stufe 1 (Cowork)","b":["🟢 Stufe 0 Notion-Ansicht","🟡 Stufe 1 Cowork (kein Code) ← jetzt","🔵 Stufe 2 Claude Code/Telegram"],"n":"Erinnerst du dich an die drei Stufen aus 4.2? Stufe null wäre einfach dein Notion aufmachen. Wir bauen jetzt Stufe eins, deinen Bot in Claude Cowork, ganz ohne Code. Und für die, die irgendwann mehr wollen, gibt's Stufe zwei mit Claude Code, der automatisch aufs Handy schickt. Für neunzig Prozent reicht Stufe eins völlig, also fangen wir da an."},
 {"t":"Stufe 1: was der Bot braucht","b":["Dein Business-Brief (aus 4.3)","Dein Notion (in Cowork verbunden)","Er liest direkt — kein Reinkopieren"],"n":"Dein Cockpit-Bot braucht genau zwei Dinge. Erstens deinen Business-Brief aus der letzten Lektion, damit er dich kennt. Und zweitens dein Notion, das du in Cowork einmal per Connector verbindest, dann liest er deine Wochenplanung direkt, du musst nichts mehr reinkopieren. Genau das machen wir jetzt zusammen."},
 {"t":"Die fertige Vorlage (Bonus)","b":["Kompletter System-Prompt fertig","Du trägst nur Namen + Kontext ein","Download direkt in der Lektion"],"n":"Damit du nicht bei null anfängst, hast du von mir eine fertige Vorlage als Bonus. Das ist der komplette System-Prompt, also die Bedienungsanleitung für den Bot, und du musst nur deinen Namen und deinen Kontext eintragen. Den Link zum Download findest du gleich hier in der Lektion."},
 {"t":"Live-Demo 1: Cowork + Notion verbinden","b":["Cowork öffnen","Notion-Connector einstecken (einloggen + freigeben)","Vorlage einfügen"],"n":"So, jetzt bau ich ihn live in Cowork, und du baust einfach mit, pausier wo du musst. Schritt eins: ich öffne Cowork, gehe in die Einstellungen und stecke den Notion-Connector ein, also einmal bei Notion einloggen und die Seiten freigeben. Dann füge ich meine fertige Vorlage ein. Du siehst, ich tippe nichts Kompliziertes, ich klicke und kopiere."},
 {"t":"Live-Demo 2: Business-Brief rein","b":["Business-Brief einfügen","Notion ist schon verbunden","Der Bot hat alles, was er braucht"],"n":"Schritt zwei: ich gebe ihm noch meinen Business-Brief, damit er weiss wer ich bin. Mein Notion ist ja jetzt verbunden, also muss ich nichts mehr reinkopieren, er schaut von selbst rein. Damit hat der Bot alles, was er braucht."},
 {"t":"Live-Demo 3: testen","b":["„Was ist heute mein Fokus?\"","→ Tagesfokus + 3 Aufgaben + Wochenblick","Wie geil ist das denn"],"n":"Und Schritt drei, der schönste: ich frag ihn „was ist heute mein Fokus\", und schau dir das an, er gibt mir meinen Tagesfokus, meine drei Hauptaufgaben und einen kurzen Wochenblick. Wie geil ist das denn. Genau das hast du gleich auch."},
 {"t":"So nutzt du ihn jeden Morgen","b":["Cowork auf → Bot fragen → loslegen","30 Sekunden statt 10 Minuten Notion","1× pro Woche neue Woche reinwerfen"],"n":"Im Alltag sieht das dann so aus: Cowork auf, Bot fragen, loslegen. Dreissig Sekunden statt zehn Minuten in Notion verlieren. Einmal pro Woche wirfst du ihm deine neue Wochenseite rein, und der Rest läuft. So fühlt sich Mama-CEO an einem normalen Morgen an."},
 {"t":"Stufe 2: der automatische Telegram-Bot","b":["Server läuft 24/7 → pusht um 6:30","Brauchst: Telegram-Bot + Hosting + etwas Code","Läuft auch wenn der Laptop zu ist","Gleicher Prompt wie eben — du schreibst NICHTS neu"],"n":"Und für die, die irgendwann Blut lecken, zeig ich dir ehrlich, wie die Profi-Variante läuft. Das ist ein Bot, der Tag und Nacht auf einem kleinen Server liegt und mir jeden Morgen um halb sieben von selbst eine Telegram-Nachricht schickt. Dafür braucht es drei Dinge: du meldest deinen Bot kostenlos bei Telegram an, du brauchst einen Ort wo er rund um die Uhr läuft, das nennt man Hosting, und du brauchst ein bisschen Code, der die Verbindung macht. Und ganz wichtig, damit dir nicht graut: deinen Bot-Prompt von eben nimmst du eins zu eins mit, du schreibst nichts Neues. Der Prompt ist das Rezept und bleibt gleich, Claude Code baut nur die Hülle drumherum, also Telegram und den automatischen Versand. Cowork ist der Herd zu Hause, Claude Code ist die Lieferung an die Haustür."},
 {"t":"Stufe 2: Kosten + Claude Code führt dich durch","b":["Hosting ~5 CHF/Mt + ggf. KI-Kosten","Claude Code schreibt den Code für dich","KEIN Muss — Stufe 1 trägt weit"],"n":"Und jetzt das Wichtige, weil ich ehrlich bleibe: dieses Hosting kostet etwa fünf Franken im Monat, dazu kommen je nach Nutzung ein paar Franken KI-Kosten. Den Code musst du nicht selbst schreiben, das macht Claude Code für dich, und es führt dich Schritt für Schritt durch alles. Aber sei dir bewusst: das ist mehr Aufwand und kostet etwas, und es ist ausdrücklich kein Muss. Stufe eins trägt dich weit, also fang da an."},
 {"t":"Arbeitsblatt + nächste Lektion","b":["📋 Cockpit-Bot bauen + 3 Tage testen","Nächste Lektion: 4.5 MASTERY — Haushalts-Helfer"],"n":"Im Arbeitsblatt baust du jetzt deinen Cockpit-Bot mit der Vorlage und testest ihn drei Tage lang jeden Morgen. In der nächsten Lektion, der MASTERY, bauen wir seinen Zwilling fürs Zuhause: deinen Haushalts-Helfer. Bis gleich."},
]

L45 = [
 {"t":"MASTERY · Dein Haushalts-Helfer-Bot","sub":"Der Mental Load raus aus deinem Kopf","n":"Willkommen zur MASTERY von Säule 4, und die ist mir besonders wichtig, weil sie zeigt, dass Mama-CEO nicht am Schreibtisch aufhört. Wir bauen jetzt den Zwilling deines Cockpit-Bots, nur diesmal für dein Zuhause, nämlich deinen Haushalts-Helfer."},
 {"t":"Der unsichtbare Mental Load","b":["Wäsche, Zahnarzt, Geschenk, Kleider, Elternabend","Liegt unsichtbar im Kopf","Lässt dich auch abends nicht los"],"n":"Mal ganz ehrlich: was wir Mamas im Kopf mit uns rumtragen, ist riesig. Wann muss die Wäsche, wer hat wann Zahnarzt, wann muss ich die Frühlingskleider raussuchen, das Geburtstagsgeschenk, der Elternabend. Das Schlimme ist, dass dieser Mental Load unsichtbar in deinem Kopf liegt und dich auch abends nicht loslässt. Genau den holen wir jetzt raus."},
 {"t":"Der Zwilling des Cockpit-Bots","b":["Gleiche Mechanik wie der Cockpit-Bot","Liest eine Liste → sagt was dran ist","Cockpit = Business, Haushalt = Zuhause"],"n":"Das Schöne ist: du kannst schon alles, was du dafür brauchst. Der Haushalts-Helfer funktioniert genau wie dein Cockpit-Bot, er liest eine Liste und sagt dir morgens, was dran ist. Der Cockpit-Bot macht das fürs Business, der Haushalts-Helfer fürs Zuhause. Zwei Zwillinge, einmal Arbeit, einmal Familie."},
 {"t":"Was er für dich tut","b":["„Montag — heute: Wäsche machen","Kind hat Zahnarzt 14 Uhr","Kind muss früher in die Schule (Wald)\""],"n":"So sieht das dann aus: du fragst morgens „was ist heute zu Hause dran\", und er sagt dir zum Beispiel: Montag, heute Wäsche machen, das Kind hat um zwei Zahnarzt, und das andere muss früher los, weil Waldtag. Alles, was sonst in deinem Kopf herumschwirrt, kommt jetzt einmal am Morgen klar auf den Tisch."},
 {"t":"Die Quelle: dein Brain Dump aus Säule 2","b":["Hütchen-Inventar = die Quelle","Alle „muss ich noch\"-Sachen","→ in eine Liste"],"n":"Und jetzt zahlt sich Säule 2 aus. Erinnerst du dich an dein Hütchen-Inventar, deinen grossen Brain Dump, wo alles rausgekommen ist, was du mit dir rumträgst? Genau das ist die Quelle. Diese ganzen muss-ich-noch-Sachen bringen wir jetzt in eine Liste, aus der dein Bot dich erinnert."},
 {"t":"Zwei Sorten Aufgaben + Termine","b":["🔁 wiederkehrend (Staubsaugen montags)","📅 datiert (Kleider raus letzter Fr im März)","👨‍👩‍👧 Familien-Termine (Zahnarzt, Waldtag)"],"n":"Deine Haushalts-Sachen sind von zwei Sorten, plus die Termine. Es gibt das Wiederkehrende, das immer am gleichen Tag dran ist, zum Beispiel Staubsaugen montags. Es gibt das Datierte, das einmal an einem bestimmten Tag kommt, zum Beispiel Frühlingskleider raussuchen am letzten Freitag im März. Und es gibt die Familien-Termine wie Zahnarzt oder Waldtag. Alle drei kommen in deine Liste."},
 {"t":"Wo die Liste lebt: Notion","b":["Eigene Haushalts-Liste in Notion","Felder: Rhythmus, Wochentag, festes Datum","Beide Bots schauen da rein"],"n":"Und wo lebt diese Liste? In Notion, genau dort, wo schon dein Business-Brain liegt. Du legst dir eine eigene Haushalts-Liste an, mit Feldern für den Rhythmus, also montags oder monatlich, und für feste Daten. So wie dein Business in Notion ein Zuhause hat, kriegt jetzt auch dein Haushalt eins, und beide Bots schauen da rein."},
 {"t":"Live-Demo 1: Haushalts-Liste füllen","b":["Notion-Haushalts-Liste öffnen","Sachen aus dem Brain Dump eintragen","Rhythmus + festes Datum dazu"],"n":"So, jetzt bauen wir live, und du baust mit, pausier wo du musst. Schritt eins: ich öffne meine Haushalts-Liste in Notion und trage ein paar Sachen aus meinem Brain Dump ein. Wäsche, Rhythmus wöchentlich. Bäder putzen, mittwochs. Kleider wechseln, festes Datum Ende März. Du füllst deine eigene Liste mit dem, was bei dir ansteht."},
 {"t":"Live-Demo 2: Bot bauen (in Cowork)","b":["Haushalts-Helfer-Vorlage in Cowork einfügen","Notion ist schon verbunden (aus 4.4)","Gleiche Mechanik, andere Liste"],"n":"Schritt zwei: ich nehme die fertige Haushalts-Helfer-Vorlage, die du als Bonus hast, und setze sie in Cowork ein. Mein Notion ist von Lektion 4.4 schon verbunden, also liest er die Haushalts-Liste direkt. Das ist exakt dieselbe Mechanik wie beim Cockpit-Bot, nur eine andere Liste und eine andere Vorlage."},
 {"t":"Live-Demo 3: testen","b":["„Was ist heute zu Hause dran?\"","→ Haushalt + Termine der Kinder","Vorher im Kopf, jetzt sagt's der Helfer"],"n":"Schritt drei, der schöne: ich frag ihn „was ist heute zu Hause dran\", und schau, er gibt mir meine Haushalts-Sachen für heute plus die Termine der Kinder. Wie geil ist das denn, das musste ich vorher alles selbst im Kopf haben, und jetzt sagt's mir mein Helfer."},
 {"t":"So nutzt du beide Bots morgens","b":["Cockpit: was im Business dran ist","Haushalts-Helfer: was zu Hause dran ist","Ganzer Tag in 2 Minuten klar"],"n":"Im Alltag fragst du morgens beide kurz: den Cockpit-Bot, was im Business dran ist, und den Haushalts-Helfer, was zu Hause dran ist. In zwei Minuten hast du deinen ganzen Tag klar vor dir, Business und Familie, ohne dass irgendwas nur noch in deinem Kopf liegt. Das ist für mich gelebter Mama-CEO."},
 {"t":"Stufe 2: auch hier automatisch möglich","b":["Wie beim Cockpit: Telegram-Push","Via Claude Code + Hosting","Kein Muss — fragen reicht"],"n":"Und ganz kurz, weil's die gleiche Logik ist wie beim Cockpit-Bot: auch deinen Haushalts-Helfer kannst du später als automatischen Telegram-Bot bauen, der dir morgens von selbst schreibt. Das braucht wieder etwas Technik und Hosting, ist also Stufe zwei und kein Muss. Für den Anfang reicht es völlig, dass du ihn fragst."},
 {"t":"Extra-Bonus: dein Kochassistent","b":["🍳 Wochenplan · Spontan-Koch · Einkaufsliste","Gleiche Mechanik, Familie als Kontext","Vorlage liegt bei"],"n":"Und weil das was-koch-ich-heute auch so ein Klassiker ist, hab ich dir noch einen Extra-Bonus dazugepackt: eine Vorlage für einen Kochassistenten. Der läuft nach genau demselben Prinzip, du gibst ihm deine Familie als Kontext, und er macht dir Wochenpläne, Spontan-Ideen und Einkaufslisten. Bau ihn dir, wann immer du Lust hast."},
 {"t":"Arbeitsblatt + nächste Lektion","b":["📋 Haushalts-Liste + Helfer bauen + testen","Nächste Lektion: 4.6 — Mensch vs. KI"],"n":"Im Arbeitsblatt bringst du jetzt deine Brain-Dump-Sachen in deine Haushalts-Liste und baust deinen Haushalts-Helfer mit der Vorlage. In der letzten Lektion von Säule 4 ziehen wir dann die Linie: was macht die KI, und was bleibt immer deins. Bis gleich."},
]

L46 = [
 {"t":"KI-Wochenplan","sub":"Mensch vs. Maschine","n":"Letzte Lektion von Säule 4, und hier ziehen wir die wichtigste Linie überhaupt: was macht ab jetzt die KI, und was bleibt immer deins. Denn beides falsch zu machen kostet dich, und ich zeig dir wie du's richtig aufteilst."},
 {"t":"Du hast jetzt 2 Mitarbeiter","b":["🌅 Cockpit-Bot · 🏠 Haushalts-Helfer","Zwei Mitarbeiter, die nie müde werden","Wie teilst du die Arbeit?"],"n":"Schau, was du in den letzten Lektionen aufgebaut hast: einen Cockpit-Bot fürs Business und einen Haushalts-Helfer fürs Zuhause. Zwei Mitarbeiter, die nie müde werden. Die Frage ist jetzt nur noch, welche Arbeit du ihnen gibst, ohne dass du dich selbst überflüssig machst."},
 {"t":"Die zwei Fehler","b":["Alles abgeben → deine Stimme weg","Nichts abgeben → zurück ins Hamsterrad","Die Kunst liegt in der Mitte"],"n":"Es gibt zwei Fehler. Der eine ist, alles an die KI abzugeben, dann klingt dein Business austauschbar und deine Kundinnen spüren, dass keine echte Frau mehr dahintersteht. Der andere ist, aus Stolz oder Angst gar nichts abzugeben, dann bleibst du im Hamsterrad aus Säule 2. Die Kunst liegt in der Mitte."},
 {"t":"Die einfache Regel","b":["KI: Wiederholbares + Vorbereitendes","Du: Beziehung + Entscheidung + Stimme","KI bereitet vor, du entscheidest"],"n":"Und die Regel dafür ist erstaunlich einfach. Die KI macht alles, was wiederholbar oder vorbereitend ist, also Recherche, erste Entwürfe, Pläne, Listen, Zusammenfassungen. Und du machst alles, wo es um Beziehung, Entscheidung und deine Stimme geht. Die KI bereitet vor, du entscheidest und gibst den letzten Schliff."},
 {"t":"Was zur KI darf","b":["Recherche, erste Entwürfe","Wochen- + Essenspläne, Listen","Zusammenfassungen, Morgenbriefing"],"n":"Konkret darf zur KI: Recherche, erste Entwürfe von Texten, deine Wochenplanung und deine Essensplanung, Listen aller Art, Zusammenfassungen von langen Sachen, und natürlich dein Morgenbriefing. All das, was Zeit frisst, aber nicht zwingend dein Herzblut braucht."},
 {"t":"Was deins bleibt: die 5 CEO-Aufgaben","b":["Vision & Strategie · Entscheidungen","Brand & Stimme · Beziehungen","Reflexion & Zahlen"],"n":"Und was bleibt immer deins? Genau die fünf CEO-Aufgaben aus Säule 1. Deine Vision und Strategie, deine Entscheidungen, deine Brand und deine Stimme, deine Beziehungen zu den Kundinnen, und deine Reflexion. Das ist der Kern, warum dich Menschen buchen, und den gibst du nie an eine Maschine ab."},
 {"t":"Dein KI-Wochenplan","b":["Jede Aufgabe stempeln: 🤖 KI / 🙋 ich","Morgenbriefing → KI, Kundengespräch → du","Erster Entwurf → KI, finale Stimme → du"],"n":"Jetzt machen wir's konkret für deine Woche. Du nimmst deine Aufgaben und stempelst hinter jede entweder KI oder ich. Das Morgenbriefing kriegt KI, das Kundengespräch kriegst du, der erste Caption-Entwurf KI, die finale Stimme du. So siehst du auf einen Blick, wo dir Stunden geschenkt werden."},
 {"t":"Vorschau Säule 5","b":["Die Linie → volle Mama-CEO-Matrix","+ System (ohne KI) + ganz raus","Das Fundament legst du hier"],"n":"Und das hier ist erst der Anfang, denn in Säule 5 bauen wir diese Linie zur vollen Mama-CEO-Matrix aus. Da kommen noch zwei Felder dazu: was kann ein System ohne KI übernehmen, und was darf einfach ganz weg. Aber das Fundament, Mensch gegen Maschine, das legst du hier."},
 {"t":"Arbeitsblatt + Säule 4 komplett","b":["📋 Dein KI-Wochenplan","Säule 4 abgeschlossen — Live-Call 3","2 Bots laufen, die Linie ist klar"],"n":"Im Arbeitsblatt baust du jetzt deinen KI-Wochenplan und stempelst jede Aufgabe. Damit hast du Säule 4 abgeschlossen: zwei KI-Mitarbeiter laufen, und du weisst genau wo Mensch aufhört und Maschine anfängt. Wir sehen uns im Live-Call 3, der Bot-Bau-Werkstatt. Ich freu mich drauf."},
]

build("01-lektion-4-1.pptx", "Mama-CEO · Säule 4 · Lektion 4.1", L41)
build("02-lektion-4-2.pptx", "Mama-CEO · Säule 4 · Lektion 4.2", L42)
build("03-lektion-4-3.pptx", "Mama-CEO · Säule 4 · Lektion 4.3", L43)
build("04-lektion-4-4.pptx", "Mama-CEO · Säule 4 · Lektion 4.4", L44)
build("05-lektion-4-5.pptx", "Mama-CEO · Säule 4 · Lektion 4.5", L45)
build("06-lektion-4-6.pptx", "Mama-CEO · Säule 4 · Lektion 4.6", L46)
print("ALL DONE")
