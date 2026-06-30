# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT = r"C:\Users\pnulm\Desktop\Mein Business\outputs\produkte\mama-ceo\03-praesentationen\saeule-5"
os.makedirs(OUT, exist_ok=True)

ORANGE = RGBColor(0xDC, 0x82, 0x2E)
NAVY   = RGBColor(0x1A, 0x3A, 0x4A)
BODY   = RGBColor(0x2C, 0x3E, 0x50)
PETROL = RGBColor(0x12, 0x82, 0x8C)
MUTED  = RGBColor(0x7A, 0x8A, 0x95)
CREME  = RGBColor(0xF1, 0xEC, 0xDD)
EYEBROW = "SÄULE 5 — BUSINESS SKALIEREN"

def set_bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = CREME

def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True; return tb, tf
def run(p, text, font="Calibri", size=16, bold=False, color=BODY):
    r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; return r
def add_title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(4.2), Inches(3.5), Inches(3.5)); c.fill.solid(); c.fill.fore_color.rgb = ORANGE; c.line.fill.background()
    _, tf = textbox(s, 0.6, 0.6, 8.8, 0.4); run(tf.paragraphs[0], EYEBROW, "Calibri", 11, True, ORANGE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05), Inches(0.6), Inches(0.05)); bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    _, tf = textbox(s, 0.6, 1.35, 8.6, 1.7); run(tf.paragraphs[0], title, "Georgia", 44, True, NAVY)
    _, tf = textbox(s, 0.6, 3.15, 8.4, 1.4); run(tf.paragraphs[0], subtitle, "Georgia", 18, False, BODY)
    _, tf = textbox(s, 0.6, 5.0, 4.0, 0.4); run(tf.paragraphs[0], "Patricia Ulmann", "Calibri", 12, True, PETROL)
    _, tf = textbox(s, 5.4, 5.0, 4.0, 0.4); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; run(p, "mumlifebalance.ch", "Calibri", 11, False, MUTED)
    return s
def add_content_slide(prs, title, bullets, tag):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    _, tf = textbox(s, 0.6, 0.5, 8.8, 1.0); run(tf.paragraphs[0], title, "Georgia", 30, True, NAVY)
    _, tf = textbox(s, 0.6, 1.55, 8.8, 3.45); first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False; p.space_after = Pt(10)
        run(p, "✓  ", "Calibri", 16, True, PETROL); run(p, b, "Calibri", 16, False, BODY)
    _, tf = textbox(s, 0.5, 5.2, 6.0, 0.3); run(tf.paragraphs[0], tag, "Calibri", 9, False, MUTED)
    _, tf = textbox(s, 7.0, 5.2, 2.5, 0.3); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; run(p, "Mum Life Balance", "Calibri", 9, False, MUTED)
    return s
def build(filename, tag, slides):
    prs = Presentation(); prs.slide_width = Inches(10); prs.slide_height = Inches(5.625)
    for i, sl in enumerate(slides):
        s = add_title_slide(prs, sl["t"], sl.get("sub","")) if i==0 else add_content_slide(prs, sl["t"], sl["b"], tag)
        if sl.get("n"): s.notes_slide.notes_text_frame.text = sl["n"]
    prs.save(os.path.join(OUT, filename)); print("saved", filename, len(slides), "Folien")

L51=[
 {"t":"Die Mama-CEO-Matrix","sub":"Was bleibt deins, was gibst du ab","n":"Willkommen in der letzten Säule. Wir ziehen jetzt alles zusammen, was du in den letzten Wochen gebaut hast, und geben deinem ganzen Business eine Ordnung. Das Werkzeug dafür heisst Mama-CEO-Matrix, und es ist dein Kompass fürs Skalieren."},
 {"t":"Wo du stehst","b":["✓ Rhythmus (Säule 1)","✓ Mindset (Säule 2)","✓ Notion-Brain (Säule 3)","✓ 2 KI-Mitarbeiter (Säule 4)"],"n":"Schau kurz, was du schon hast: deinen Rhythmus, dein Mindset, dein Notion-Brain und zwei KI-Mitarbeiter. Das ist enorm viel. Jetzt fehlt nur noch eins, nämlich für jede einzelne Aufgabe zu wissen, wo sie hingehört. Genau das macht die Matrix."},
 {"t":"Von der Linie zur Matrix","b":["Säule 4: ich / KI (zwei Seiten)","+ ⚙️ System (läuft ohne dich)","+ 🗑 raus (darf aufhören)"],"n":"In Säule 4 hatten wir zwei Seiten: was macht die KI, was machst du. Das war der Anfang. Jetzt kommen zwei Felder dazu, die das Bild komplett machen: Sachen, die ganz ohne dich als System laufen, und Sachen, die einfach weg dürfen. Aus der Linie wird ein Vierer-Raster."},
 {"t":"Die 4 Felder","b":["🙋 ich — nur du","🤖 KI — Bots bereiten vor","⚙️ System — läuft mechanisch ohne dich","🗑 raus — darf aufhören"],"n":"Das ist die Matrix: vier Felder. Ich, also nur du. KI, deine Bots bereiten vor. System, läuft mechanisch ohne dich. Und raus, das darf einfach aufhören. Jede Aufgabe deines Business landet in genau einem dieser Felder, und das gibt dir eine Klarheit, die du vorher nie hattest."},
 {"t":"🙋 ICH = nur du","b":["Die 5 CEO-Aufgaben (aus Säule 1)","1:1 mit deinen Kundinnen","Deine Story, deine Stimme"],"n":"Fangen wir beim wichtigsten Feld an: ich. Da gehören deine fünf CEO-Aufgaben rein, die du aus Säule 1 kennst, plus die echten Beziehungen, die Gespräche mit deinen Kundinnen. Das ist der Kern, warum dich Menschen buchen, und der bleibt immer deins."},
 {"t":"🤖 KI vs. ⚙️ System","b":["KI denkt mit: Auftrag → Output","System läuft mechanisch: Vorlage, Automatik, Ämtli","Faustregel: jedes Mal urteilen = KI · immer gleich = System"],"n":"Jetzt der feine, aber wichtige Unterschied. Die KI denkt mit: du gibst einen Auftrag, sie liefert was zurück, zum Beispiel dein Morgenbriefing oder ein erster Entwurf. Ein System läuft mechanisch, ohne Denken: eine Vorlage, eine Automatisierung, ein festes Ämtli in der Familie. Faustregel: musst du jedes Mal urteilen, ist es KI. Läuft es immer gleich, ist es System."},
 {"t":"🗑 RAUS — das mutigste Feld","b":["Was eigentlich niemand braucht","Aus dem 4-Filter „X\" (Säule 2)","Weglassen ist die ehrlichste Form von Skalieren"],"n":"Und dann das Feld, das die meisten überspringen, weil es wehtut: raus. Das sind die Sachen, die eigentlich niemand braucht und die du nur aus Gewohnheit machst. Erinnerst du dich an den 4-Filter aus Säule 2, das X? Genau die kommen hierher. Etwas wegzulassen ist die ehrlichste Form von Skalieren."},
 {"t":"Arbeitsblatt + nächste Lektion","b":["📋 Zu jedem Feld 1 Beispiel notieren","Nächste Lektion: 5.2 MASTERY — 25 Tasks sortieren"],"n":"Im Arbeitsblatt schreibst du zu jedem Feld schon mal ein Beispiel aus deinem Alltag. In der nächsten Lektion, der MASTERY, machen wir's dann richtig: du nimmst fünfundzwanzig echte Aufgaben und sortierst jede in ihr Feld. Bis gleich."},
]
L52=[
 {"t":"MASTERY · 25 Tasks sortieren","sub":"Die Matrix füllen — ich / KI / System / raus","n":"Das ist die Lektion, auf die alles hinausläuft. Wir füllen jetzt deine Matrix mit echten Aufgaben, und ich verspreche dir, am Ende schaust du anders auf dein Business. Nimm dir Zeit und ein ungestörtes halbes Stündchen."},
 {"t":"Warum genau 25","b":["5 ist zu wenig, um ein Muster zu sehen","Alles auf einmal überfordert","25 macht alles sichtbar"],"n":"Warum fünfundzwanzig? Weil fünf zu wenig sind, um ein Muster zu sehen, und alles auf einmal überfordert. Fünfundzwanzig ist genau die Menge, bei der sichtbar wird, wie dein Business wirklich verteilt ist. Du wirst überrascht sein."},
 {"t":"Schritt 1: 25 Aufgaben sammeln","b":["Aus deinem Hütchen-Inventar (Säule 2)","+ dein Business-Alltag (Content, DMs, Buchhaltung …)","Einfach runterschreiben, was du in einer Woche tust"],"n":"Schritt eins: sammle fünfundzwanzig echte Aufgaben. Nimm dein Hütchen-Inventar aus Säule 2 und ergänze deinen Business-Alltag, also Content, DMs, Buchhaltung, Planung, was auch immer dich beschäftigt. Schreib einfach runter, was du in einer typischen Woche so tust."},
 {"t":"Schritt 2: jede in EIN Feld","b":["Genau ein Feld pro Aufgabe","Nicht „kommt drauf an\" — entscheide dich","Erstes ehrliches Bauchgefühl (korrigieren geht später)"],"n":"Schritt zwei: jede Aufgabe kommt in genau ein Feld. Nicht zwei, nicht kommt drauf an, entscheide dich. Und keine Sorge, du darfst später korrigieren. Es geht erstmal ums erste, ehrliche Bauchgefühl."},
 {"t":"Hilfsfrage 🙋 ICH","b":["„Würde es meiner Marke / Beziehung schaden,","wenn das jemand anders macht?\"","Ja → bleibt bei dir · Nein → weg vom Tisch"],"n":"Damit du nicht raten musst, gebe ich dir pro Feld eine Hilfsfrage. Für ich frag dich: würde es meiner Marke oder einer Beziehung schaden, wenn das jemand anders macht? Wenn ja, gehört es zu dir. Wenn nein, darf es weg von deinem Tisch."},
 {"t":"Hilfsfrage 🤖 KI","b":["„Vorbereitend oder wiederholbar,","braucht aber Urteil?\"","Recherche, Entwurf, Plan → ab zur KI"],"n":"Für KI frag dich: ist das vorbereitend oder wiederholbar, braucht aber trotzdem ein bisschen Urteil? Recherche, ein erster Entwurf, ein Plan. Wenn ja, ab zur KI. Deine Bots aus Säule 4 sind genau dafür da."},
 {"t":"Hilfsfrage ⚙️ SYSTEM","b":["„Kann das mechanisch laufen","oder jemand fix übernehmen?\"","Vorlage, Automatik, Ämtli → einmal bauen, nie wieder anfassen"],"n":"Für System frag dich: kann das immer gleich ablaufen, ohne dass jemand denken muss, oder kann es jemand anders fix übernehmen? Eine Vorlage, eine Automatik, ein Ämtli für die Kinder. Wenn ja, mach ein System draus, dann musst du es nie wieder anfassen."},
 {"t":"Hilfsfrage 🗑 RAUS","b":["„Braucht das wirklich jemand?","Was passiert, wenn ich's lasse?\"","Oft: nichts passiert → darf weg"],"n":"Und für raus die ehrlichste Frage: braucht das wirklich jemand, und was passiert, wenn ich es einfach lasse? Oft ist die Antwort: nichts passiert. Genau das darf weg, und das ist kein Versagen, das ist Führung."},
 {"t":"Live-Demo: meine Matrix","b":["Patricias eigene Tasks in den 4 Feldern","Feld „ich\" ist erstaunlich klein","Viel läuft über KI + Systeme"],"n":"Ich zeig dir jetzt meine eigene Matrix, damit du ein Gefühl kriegst, wie das aussieht. Du siehst, bei mir ist das Feld ich erstaunlich klein, und ganz viel läuft über KI und Systeme. Genau dahin willst du auch."},
 {"t":"Der Aha: das meiste ist NICHT „ich\"","b":["Feld „ich\" klein, Rest gross","Die Arbeit, die dich auffrisst, ist gar nicht deine","Das ist deine Erlaubnis loszulassen"],"n":"Und das ist der grosse Aha-Moment dieser Lektion: das Feld ich ist viel kleiner, als du denkst. Die meiste Arbeit, die dich auffrisst, gehört eigentlich gar nicht zu dir. Das ist die Erlaubnis loszulassen, schwarz auf weiss."},
 {"t":"Dein nächster Schritt","b":["Aus KI/System die 2-3 nervigsten rauspicken","Einen Bot mehr · eine Automatik · ein Ämtli abgeben","Schritt für Schritt wird dein Tisch leerer"],"n":"Damit das nicht nur ein schönes Bild bleibt: such dir aus deinem KI- und System-Feld die zwei, drei Sachen raus, die dich am meisten nerven, und richte sie als Erstes ein. Ein Bot mehr, eine Automatik, ein Ämtli abgegeben. Schritt für Schritt wird dein Tisch leerer."},
 {"t":"Arbeitsblatt + nächste Lektion","b":["📋 Deine 25-Task-Matrix füllen (in Notion)","Nächste Lektion: 5.3 — was deins bleibt"],"n":"Im Arbeitsblatt füllst du jetzt deine eigene 25-Task-Matrix, am besten direkt in deinem Notion-Template. In der nächsten Lektion schauen wir genauer auf das Feld ich: welche Beziehungen du beim Skalieren auf keinen Fall verlierst. Bis gleich."},
]
L53=[
 {"t":"Dein innerer Kreis","sub":"Was DEINS bleibt — Beziehungen schützen","n":"In der letzten Lektion hast du gesehen, wie viel du abgeben kannst. Jetzt reden wir über das Gegenteil: über das, was du auf keinen Fall abgibst. Denn beim Skalieren kann man sich auch kaputt-automatisieren, und davor will ich dich bewahren."},
 {"t":"Die Gefahr beim Skalieren","b":["Effizienz ↔ Nähe (Waage)","Zu viel Automatik fühlt sich kühl an","Skalieren nie auf Kosten der Beziehung"],"n":"Die Gefahr ist verlockend: du wirst so gut im Automatisieren, dass du irgendwann alles an Bots und Systeme gibst, auch das Menschliche. Und dann wundern sich deine Kundinnen, warum sich alles so kühl anfühlt. Effizienz ist super, aber nicht auf Kosten der Nähe."},
 {"t":"Der innere Kreis","b":["Die Beziehungen, die immer von DIR kommen","Egal wie gross dein Business wird","Alles andere darf über KI/System laufen"],"n":"Darum definierst du jetzt deinen inneren Kreis. Das sind die Menschen und Beziehungen, die immer direkt von dir kommen, egal wie gross dein Business wird. Alles andere darf über KI und System laufen, aber dieser Kreis bleibt heilig."},
 {"t":"Bestandskundinnen = Gold","b":["Die, die schon ja gesagt haben","Halten ist günstiger als neu gewinnen","Diese Betreuung nie an die KI"],"n":"Ganz vorne in diesem Kreis: deine Bestandskundinnen. Die haben schon ja zu dir gesagt, die kennen dich, die empfehlen dich weiter. Eine neue Kundin zu gewinnen ist viel teurer, als eine bestehende zu halten. Diese Gespräche, diese Betreuung, die gibst du nie an eine Maschine."},
 {"t":"Sparring + Mentorin","b":["Dein Frauen-Team / Austausch auf Augenhöhe","Die Person, von der du lernst","Tragen dich an schwierigen Tagen"],"n":"Dazu gehören auch die Frauen, mit denen du dich austauschst, dein Sparring, und die Person, von der du lernst. Diese Beziehungen tragen dich an schwierigen Tagen, und sie brauchen echte Zeit von dir, keine automatische Nachricht."},
 {"t":"Familie zuerst","b":["Der innerste Kreis (aus Säule 1/2)","Kein System ist es wert, dass Familie hintenrunterfällt","Dafür skalierst du ja überhaupt"],"n":"Und der allerinnerste Kreis ist deine Familie. Das hast du schon in Säule 1 und 2 verankert, und es bleibt der Massstab: kein System der Welt ist es wert, dass die Familie hintenrunterfällt. Deshalb skalierst du ja überhaupt, um mehr von dieser Zeit zu haben."},
 {"t":"Wie du den Kreis schützt","b":["Feste Zeit reservieren (im Wochenrhythmus)","Was einen festen Platz hat, fällt nicht weg","Was „wenn Zeit bleibt\" ist, verschwindet zuerst"],"n":"Schützen heisst nicht hoffen, sondern planen. Reservier deinem inneren Kreis feste Zeit in deinem Wochenrhythmus, genau wie du es mit deinen Power-Slots machst. Was einen festen Platz hat, fällt nicht weg. Was nur wenn-Zeit-bleibt ist, verschwindet als Erstes."},
 {"t":"Arbeitsblatt + nächste Lektion","b":["📋 Definier deinen inneren Kreis + schütz ihn","Nächste Lektion: 5.4 — dein 90-Tage-Plan"],"n":"Im Arbeitsblatt schreibst du deinen inneren Kreis auf und legst fest, wie du ihm feste Zeit gibst. In der letzten Lektion bauen wir dann deinen 90-Tage-Plan, damit du genau weisst, wie's nach dem Programm weitergeht. Bis gleich."},
]
L54=[
 {"t":"Dein 90-Tage-Plan","sub":"Was kommt als Nächstes, was machst du NICHT","n":"Das ist sie, die letzte Lektion des ganzen Programms. Wir machen jetzt aus all deiner Arbeit einen konkreten Plan für die nächsten neunzig Tage, damit du nicht mit einem schönen Gefühl, sondern mit klaren Schritten rausgehst."},
 {"t":"Von der Methode zum konkreten Plan","b":["Säule 3 = Methode gelernt","Jetzt: DEINE nächsten 3 Monate","Am Ende weisst du jeden Tag, was zu tun ist"],"n":"In Säule 3 hast du die Methode gelernt, wie man von oben nach unten plant. Jetzt machen wir's konkret für dich: einen klaren Plan für die nächsten drei Monate. Am Ende dieser Lektion weisst du an jedem einzelnen Tag, was zu tun ist, und musst nie wieder morgens überlegen, wo du anfängst."},
 {"t":"Pro Monat EIN Zentrum","b":["1 Thema · 1 0€-Angebot · 1 Produkt","Das Monats-Trio aus Säule 3","Eines pro Monat, richtig ausgespielt"],"n":"Das Geheimnis ist Fokus: pro Monat gibt es EIN Zentrum. Ein Thema, über das du redest, ein 0-Euro-Angebot, mit dem du Leute reinholst, und ein Produkt, das du verkaufst. Genau das Monats-Trio aus Säule 3. Nicht zehn Themen gleichzeitig, eines pro Monat, das du richtig ausspielst."},
 {"t":"Beispiel: 3 Monate auf einen Blick","b":["Tabelle: Monat 1/2/3","je Zeile: Thema · 0€-Angebot · Produkt","Mein Beispiel — dann füllst du deins"],"n":"Ich zeig dir meine eigenen nächsten drei Monate als Beispiel, damit du siehst, wie schlicht das aussieht. Eine Zeile pro Monat: Thema, 0-Euro-Angebot, Produkt. Mehr braucht das Grundgerüst nicht. Du füllst gleich deine eigenen drei Zeilen."},
 {"t":"Rückwärts rechnen","b":["Vom Monatsziel → Wochen → Aufgaben","Launch rückwärts denken","Grosses Ziel → viele kleine Aufgaben"],"n":"Und jetzt der wichtigste Denk-Dreh: du planst rückwärts. Du nimmst dein Monatsziel, zum Beispiel diesen Launch am Monatsende, und rechnest rückwärts, was vorher passieren muss. Welche Mails, welche Story-Tage, welcher Content, welches Freebie zuerst. So zerfällt ein grosses Ziel in lauter kleine, machbare Aufgaben."},
 {"t":"Ab in die Aufgabenliste","b":["Aufgaben → Notion (aus Säule 3)","Cockpit-Bot serviert sie dir täglich","Du musst den Plan nie mehr suchen"],"n":"Diese Aufgaben schreibst du nicht auf einen Zettel, der verschwindet, sondern direkt in deine Notion-Aufgabenliste aus Säule 3. Und jetzt zahlt sich alles aus: dein Cockpit-Bot aus Säule 4 serviert dir jeden Morgen genau die Aufgaben, die heute dran sind. Du musst nie wieder den Plan suchen, er kommt zu dir."},
 {"t":"Wer macht was? Die Matrix entscheidet","b":["Jede Aufgabe: 🙋 ich / 🤖 KI / ⚙️ System","Die Hälfte landet gar nicht auf deinem Tisch","Plan = klar UND leicht"],"n":"Und bei jeder Aufgabe stellst du die Matrix-Frage aus dieser Säule: mach ich das selbst, oder kann es die KI vorbereiten, oder läuft es über ein System? So ist dein Plan nicht nur klar, sondern auch leicht, weil die Hälfte gar nicht mehr auf deinem Tisch landet."},
 {"t":"Die Nicht-Liste","b":["Was du bewusst NICHT machst","Kein neues Tool, keine neue Plattform, kein 5. Thema","Schützt deinen Fokus mehr als jede To-do-Liste"],"n":"Und mein Lieblingsteil, der fast nie gemacht wird: schreib auf, was du in diesen drei Monaten bewusst NICHT machst. Kein neues Tool, keine neue Plattform, kein fünftes Thema. Diese Nicht-Liste schützt deinen Fokus mehr als jede To-do-Liste, weil sie dich vor dir selbst bewahrt."},
 {"t":"Fertig besser als perfekt + wie's weitergeht","b":["Dein Plan darf Version 1 sein","Die Community bleibt offen","Cockpit-Bot begleitet dich jeden Morgen"],"n":"Zwei Sätze zum Schluss. Erstens: fertig ist besser als perfekt, dein Plan darf Version eins sein, du verbesserst im Tun. Und zweitens: du bist nicht allein. Die Community bleibt offen, und dein Cockpit-Bot begleitet dich jeden Morgen weiter. Das Programm endet, dein System läuft weiter."},
 {"t":"Abschluss — du bist Mama-CEO","b":["Zeit · Mindset · Struktur · KI · Matrix · Plan","📋 Plan im Arbeitsblatt ausfüllen","Du führst statt zu funktionieren 💛"],"n":"Schau, was du in acht Wochen gebaut hast: Zeit, ein neues Mindset, eine Struktur in Notion, zwei KI-Mitarbeiter, eine klare Matrix und jetzt einen konkreten Plan für drei Monate. Im Arbeitsblatt füllst du diesen Plan aus. Und dann sag ich dir von Herzen: du bist jetzt Mama-CEO, nicht weil alles perfekt ist, sondern weil du führst statt funktionierst. Wir sehen uns im Live-Call vier, bring deinen Plan mit."},
]

build("01-lektion-5-1.pptx", "Mama-CEO · Säule 5 · Lektion 5.1", L51)
build("02-lektion-5-2.pptx", "Mama-CEO · Säule 5 · Lektion 5.2", L52)
build("03-lektion-5-3.pptx", "Mama-CEO · Säule 5 · Lektion 5.3", L53)
build("04-lektion-5-4.pptx", "Mama-CEO · Säule 5 · Lektion 5.4", L54)
print("ALL DONE")
